#!/usr/bin/env python3
"""Claude Code 워크북 핸드아웃 — 28장 PPTX.

handout.md 전체 내용을 28장 슬라이드로.
브랜드: 미래에셋생명 · 오렌지 #F58220 + 짙은 네이비.
"""
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

DOCS = Path(__file__).parent
OUT = DOCS / "handout.pptx"

# ── 브랜드 색상 ────────────────────────────────────────
OR = RGBColor(0xF5, 0x82, 0x20)
OR_D = RGBColor(0xCB, 0x60, 0x15)
NV = RGBColor(0x04, 0x3B, 0x72)
NV_D = RGBColor(0x06, 0x1E, 0x30)
BG = RGBColor(0x04, 0x18, 0x28)
TX = RGBColor(0xE5, 0xE8, 0xEC)
TX2 = RGBColor(0x8D, 0xA0, 0xB8)
TX3 = RGBColor(0x5A, 0x7A, 0x98)
GD = RGBColor(0x86, 0xEF, 0xAC)
GD_D = RGBColor(0x10, 0xB9, 0x81)
BAD = RGBColor(0xFC, 0xA5, 0xA5)
WN = RGBColor(0xFB, 0xBF, 0x24)
BL = RGBColor(0x60, 0xA5, 0xFA)
BLM = RGBColor(0x38, 0xBD, 0xF8)
BD = RGBColor(0x0A, 0x30, 0x50)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BLACK = RGBColor(0x00, 0x00, 0x00)

KO_FONT = "Pretendard"
MONO = "JetBrains Mono"

# ── 슬라이드 사이즈: 16:9 와이드 ───────────────────────
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

SLIDE_W = prs.slide_width
SLIDE_H = prs.slide_height
BLANK = prs.slide_layouts[6]


# ── 헬퍼 ──────────────────────────────────────────────
def add_bg(slide, color=BG):
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H)
    bg.fill.solid()
    bg.fill.fore_color.rgb = color
    bg.line.fill.background()
    return bg


def add_text(slide, x, y, w, h, text, *, size=18, color=TX, bold=False,
             font=KO_FONT, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP):
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Emu(0)
    tf.margin_top = tf.margin_bottom = Emu(0)
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return box


def add_rect(slide, x, y, w, h, *, fill=None, line=None, line_w=None, radius=False):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE,
        x, y, w, h,
    )
    if fill is not None:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    else:
        shape.fill.background()
    if line is not None:
        shape.line.color.rgb = line
        if line_w:
            shape.line.width = line_w
    else:
        shape.line.fill.background()
    return shape


def add_chip(slide, x, y, text, *, fill=OR, fg=WHITE, size=11, padx=8):
    w = Emu(int(len(text) * 110000 + padx * 18000))
    h = Inches(0.32)
    chip = add_rect(slide, x, y, w, h, fill=fill, radius=True)
    chip.line.fill.background()
    tf = chip.text_frame
    tf.margin_left = tf.margin_right = Emu(60000)
    tf.margin_top = tf.margin_bottom = Emu(0)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = text
    r.font.name = KO_FONT
    r.font.size = Pt(size)
    r.font.bold = True
    r.font.color.rgb = fg
    return chip


def add_header(slide, kicker, title):
    add_chip(slide, Inches(0.6), Inches(0.5), kicker)
    add_text(
        slide, Inches(0.6), Inches(0.95), Inches(12), Inches(0.9),
        title, size=30, bold=True, color=WHITE,
    )
    line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0.6), Inches(1.85), Inches(0.7), Emu(40000),
    )
    line.fill.solid()
    line.fill.fore_color.rgb = OR
    line.line.fill.background()
    add_text(
        slide, Inches(12.5), Inches(7.05), Inches(0.7), Inches(0.3),
        f"{len(prs.slides) + 1:02d}",
        size=10, color=TX3, font=MONO, align=PP_ALIGN.RIGHT,
    )


def add_card(slide, x, y, w, h, *, accent=OR):
    card = add_rect(slide, x, y, w, h, fill=NV_D, line=BD, line_w=Pt(0.75), radius=True)
    add_rect(slide, x, y, Inches(0.06), h, fill=accent)
    return card


def bullet(tf, text, *, size=13, color=TX, bold=False, indent=0, mono=False, first=False):
    p = tf.paragraphs[0] if first else tf.add_paragraph()
    p.alignment = PP_ALIGN.LEFT
    p.level = indent
    p.space_after = Pt(3)
    r = p.add_run()
    r.text = text
    r.font.name = MONO if mono else KO_FONT
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.color.rgb = color
    return p


def card_with_bullets(slide, x, y, w, h, accent, title, bullets_list, *, title_size=16):
    add_card(slide, x, y, w, h, accent=accent)
    add_text(slide, x + Inches(0.3), y + Inches(0.18), w - Inches(0.4), Inches(0.4),
             title, size=title_size, bold=True, color=accent)
    box = slide.shapes.add_textbox(
        x + Inches(0.3), y + Inches(0.65),
        w - Inches(0.5), h - Inches(0.8),
    )
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Emu(0)
    for i, ln in enumerate(bullets_list):
        bullet(tf, ln, size=12, color=TX, first=(i == 0))


def make_table_slide(slide, kicker, title, headers, rows, *, col_widths=None):
    """헤더 + 표를 가진 슬라이드."""
    add_header(slide, kicker, title)
    cols = len(headers)
    tbl_rows = len(rows) + 1
    if col_widths is None:
        total_w = 12.0
        col_widths = [total_w / cols] * cols
    tbl_w = sum(col_widths)
    row_h = 0.45
    tbl = slide.shapes.add_table(
        tbl_rows, cols,
        Inches(0.6), Inches(2.2),
        Inches(tbl_w), Inches(row_h * tbl_rows),
    ).table
    # 헤더 행
    for ci, h in enumerate(headers):
        tbl.columns[ci].width = Inches(col_widths[ci])
        cell = tbl.cell(0, ci)
        cell.text = h
        for p in cell.text_frame.paragraphs:
            p.font.name = KO_FONT
            p.font.size = Pt(11)
            p.font.bold = True
            p.font.color.rgb = WHITE
        cell.fill.solid()
        cell.fill.fore_color.rgb = NV
    # 데이터 행
    for ri, row in enumerate(rows):
        for ci, val in enumerate(row):
            cell = tbl.cell(ri + 1, ci)
            cell.text = val
            for p in cell.text_frame.paragraphs:
                p.font.name = KO_FONT
                p.font.size = Pt(10)
                p.font.color.rgb = TX
            cell.fill.solid()
            cell.fill.fore_color.rgb = NV_D if ri % 2 == 0 else BG


# ═══════════════════════════════════════════════════════
# Slide 1 — 표지
# ═══════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
add_bg(s)
add_rect(s, 0, 0, Inches(5), SLIDE_H, fill=OR_D)
add_text(s, Inches(0.6), Inches(0.6), Inches(4), Inches(0.5),
         "MIRAE ASSET LIFE", size=14, bold=True, color=WHITE)
add_text(s, Inches(0.6), Inches(1.0), Inches(4), Inches(0.4),
         "Vibe Coding Workshop", size=12, color=WHITE)

add_text(s, Inches(5.6), Inches(1.7), Inches(7.4), Inches(0.5),
         "WORKSHOP HANDOUT", size=12, bold=True, color=OR)
add_text(s, Inches(5.6), Inches(2.2), Inches(7.4), Inches(2.2),
         "Claude Code\n워크북", size=48, bold=True, color=WHITE)
add_text(s, Inches(5.6), Inches(4.8), Inches(7.4), Inches(0.8),
         "미래에셋생명 · AI 코딩 어시스턴트 활용 워크숍\n5시간 과정 가이드북",
         size=18, color=TX2)
add_text(s, Inches(5.6), Inches(6.6), Inches(7.4), Inches(0.4),
         "미래에셋생명  ·  2026", size=11, color=TX3, font=MONO)

# ═══════════════════════════════════════════════════════
# Slide 2 — 목차
# ═══════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
add_bg(s)
add_header(s, "CONTENTS", "목차")

toc = [
    ("Part 1", "도입 — 왜 AI 코딩 어시스턴트인가", OR,
     ["1.1 LLM이 바꾼 업무 방식", "1.2 챗봇 vs 에이전트", "1.3 실제 사례",
      "1.4 오늘의 목표", "1.5 워크숍 환경", "1.6 오늘 배울 도구"]),
    ("Part 2", "개념 설명 및 시연 — 핵심 기능", BLM,
     ["2.1~2.8 시작/프롬프트/Plan/권한/CLAUDE.md/Skill/Command/Hook",
      "2.9~2.12 컨텍스트/병렬/성공패턴/주의사항"]),
    ("Part 3", "기능 체험 — 8개의 미션", GD,
     ["미션 1~4: Plan / 권한 / CLAUDE.md / Skill",
      "미션 5~8: Command / Hook / 컨텍스트 / 노트"]),
    ("Part 4", "실습 프로젝트", OR_D,
     ["7단계 워크플로우", "완성 후 손에 남는 것"]),
    ("Appendix A", "워크숍 이후 — 집에서 혼자 쓰기", TX3,
     ["A.1~A.2 CLI/Desktop 설치", "A.3~A.5 비교/명령어/안전"]),
]
y = Inches(2.2)
for part, title, accent, items in toc:
    add_chip(s, Inches(0.6), y, part, fill=accent)
    add_text(s, Inches(2.2), y, Inches(5), Inches(0.35),
             title, size=14, bold=True, color=WHITE)
    for i, item in enumerate(items):
        add_text(s, Inches(7.5), y + Inches(i * 0.28), Inches(5.5), Inches(0.28),
                 item, size=10, color=TX2)
    y += Inches(max(len(items) * 0.28, 0.35) + 0.35)

# ═══════════════════════════════════════════════════════
# Slide 3 — 1.1 LLM이 바꾼 업무 방식
# ═══════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
add_bg(s)
add_header(s, "1.1", "LLM이 바꾼 업무 방식")

# 예전 vs 지금
card_with_bullets(
    s, Inches(0.6), Inches(2.3), Inches(5.8), Inches(4.5), TX3,
    "예전에는",
    [
        "AI에게 시키려면 코드를 직접 짜야 했음",
        "보고서 한 편: 자료조사→분류→작성→검수",
        "며칠이 걸렸음",
        "",
        "사람이 모든 것을 수동으로 처리",
    ],
)
card_with_bullets(
    s, Inches(6.7), Inches(2.3), Inches(6.2), Inches(4.5), GD,
    "지금은",
    [
        "자연어로 말하면 AI가 코드를 만들고 실행",
        "보고서·분석·자동화 결과물이 곧장 완성",
        '사람은 "무엇을, 왜, 어떻게"만 결정',
        "",
        "이것을 가능하게 하는 도구:",
        "Claude Code — Anthropic의 공식 CLI",
    ],
)

# ═══════════════════════════════════════════════════════
# Slide 4 — 1.2 챗봇 vs 에이전트
# ═══════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
add_bg(s)
add_header(s, "1.2", "챗봇과 에이전트는 어떻게 다른가요")

make_table_slide.__wrapped__ = True  # skip — we do it manually
headers = ["구분", "챗봇 (ChatGPT 등)", "에이전트 (Claude Code)"]
rows = [
    ["입력", "텍스트 질문", "텍스트 지시"],
    ["출력", "텍스트 답변", "파일 생성 · 코드 실행 · 결과물"],
    ["컴퓨터 접근", "없음", "직접 파일 · 터미널 사용"],
    ["다단계 작업", "사람이 단계마다 입력", "알아서 끝까지"],
    ["결과물 종류", "텍스트 한 덩어리", "보고서, 코드, 데이터, 이미지 등"],
]
col_w = [2.5, 4.5, 5.0]
tbl = s.shapes.add_table(
    len(rows) + 1, 3,
    Inches(0.6), Inches(2.2),
    Inches(sum(col_w)), Inches(0.45 * (len(rows) + 1)),
).table
for ci, h in enumerate(headers):
    tbl.columns[ci].width = Inches(col_w[ci])
    cell = tbl.cell(0, ci)
    cell.text = h
    for p in cell.text_frame.paragraphs:
        p.font.name = KO_FONT; p.font.size = Pt(12); p.font.bold = True; p.font.color.rgb = WHITE
    cell.fill.solid(); cell.fill.fore_color.rgb = NV
for ri, row in enumerate(rows):
    for ci, val in enumerate(row):
        cell = tbl.cell(ri + 1, ci)
        cell.text = val
        for p in cell.text_frame.paragraphs:
            p.font.name = KO_FONT; p.font.size = Pt(11); p.font.color.rgb = TX
        cell.fill.solid()
        cell.fill.fore_color.rgb = NV_D if ri % 2 == 0 else BG

add_text(s, Inches(0.6), Inches(5.3), Inches(12), Inches(0.8),
         "핵심: 챗봇은 \"말로 답합니다\", 에이전트는 \"행동합니다\"\n"
         "비유: 챗봇 = 카카오톡 상담사, 에이전트 = 잘 훈련된 신입 인턴",
         size=14, color=GD)

# ═══════════════════════════════════════════════════════
# Slide 5 — 1.3 실제 사례
# ═══════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
add_bg(s)
add_header(s, "1.3", "실제 사례 — 다른 회사들이 얻은 효과")

cases_h = ["조직 / 사례", "효과", "내용"]
cases = [
    ["Anthropic 마케팅팀", "2.5h → 30min", "고객 사례 보고서 작성 시간 단축"],
    ["ServiceNow (29,000명)", "준비 시간 95%↓", "영업 미팅 준비 자료 자동 생성"],
    ["Block (결제 기업)", "운영 효율 40%↑", "결제·상점 데이터 실시간 문서화"],
    ["DevOps 엔지니어", "5초 → 0.5초", "Redis CPU 급증 원인 5분 만에 분석"],
    ["비개발자 마케터", "30분 → 30초", "Figma 광고 시안 플러그인 코드 없이 제작"],
    ["엔터프라이즈 평균", "생산성 30~60%↑", "코드리뷰·디버깅·문서화 자동화"],
]
cases_cw = [3.5, 2.5, 6.0]
tbl = s.shapes.add_table(
    len(cases) + 1, 3,
    Inches(0.6), Inches(2.2),
    Inches(sum(cases_cw)), Inches(0.42 * (len(cases) + 1)),
).table
for ci, h in enumerate(cases_h):
    tbl.columns[ci].width = Inches(cases_cw[ci])
    cell = tbl.cell(0, ci)
    cell.text = h
    for p in cell.text_frame.paragraphs:
        p.font.name = KO_FONT; p.font.size = Pt(11); p.font.bold = True; p.font.color.rgb = WHITE
    cell.fill.solid(); cell.fill.fore_color.rgb = NV
for ri, row in enumerate(cases):
    for ci, val in enumerate(row):
        cell = tbl.cell(ri + 1, ci)
        cell.text = val
        for p in cell.text_frame.paragraphs:
            p.font.name = KO_FONT; p.font.size = Pt(10); p.font.color.rgb = TX
        cell.fill.solid()
        cell.fill.fore_color.rgb = NV_D if ri % 2 == 0 else BG

add_text(s, Inches(0.6), Inches(5.6), Inches(12), Inches(0.5),
         "AI 도입은 더 이상 \"코딩 잘하는 사람만\"의 영역이 아닙니다.\n"
         "비개발자도 한 줄 지시로 실제 결과물을 만들고 있습니다.",
         size=13, color=TX2)

# ═══════════════════════════════════════════════════════
# Slide 6 — 1.4 오늘의 목표
# ═══════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
add_bg(s)
add_header(s, "1.4", "오늘의 목표")

goals = [
    ("1", "자연어 한 줄로 보고서를 만들 수 있습니다", OR),
    ("2", "본인 부서의 문서 자동화 프로그램을 직접 설계할 수 있습니다", BLM),
    ("3", "CLAUDE.md, Skill, Command, Hook 등 핵심 기능 5가지를 활용할 수 있습니다", GD),
    ("4", "새로운 자동화가 필요할 때 혼자서도 같은 패턴으로 만들 수 있습니다", OR_D),
]
for i, (num, text, accent) in enumerate(goals):
    cy = Inches(2.4 + i * 1.2)
    add_card(s, Inches(0.6), cy, Inches(12), Inches(1.0), accent=accent)
    circle = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.85), cy + Inches(0.18), Inches(0.6), Inches(0.6))
    circle.fill.solid(); circle.fill.fore_color.rgb = accent
    circle.line.fill.background()
    add_text(s, Inches(0.85), cy + Inches(0.18), Inches(0.6), Inches(0.6),
             num, size=20, bold=True, color=WHITE, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, Inches(1.7), cy + Inches(0.25), Inches(10.5), Inches(0.5),
             text, size=15, color=TX)

add_text(s, Inches(0.6), Inches(6.9), Inches(12), Inches(0.4),
         "핵심 다섯 가지 + 한 개의 완성된 결과물, 이 두 가지를 들고 가시는 것이 목표입니다.",
         size=12, color=GD)

# ═══════════════════════════════════════════════════════
# Slide 7 — 1.5 워크숍 환경
# ═══════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
add_bg(s)
add_header(s, "1.5", "워크숍 환경 안내")

# 상단 컨트롤
add_text(s, Inches(0.6), Inches(2.1), Inches(12), Inches(0.4),
         "공통 컨트롤 (모든 슬라이드 상단)", size=14, bold=True, color=OR)
ctrls = [
    ("전체화면", "F5 키도 가능"),
    ("라이트/다크", "모드 전환"),
    ("확대·축소", "7단계"),
    ("이전/다음", "화살표 키 가능"),
    ("강의 모드", "자동 슬라이드 진행"),
]
for i, (name, desc) in enumerate(ctrls):
    cx = Inches(0.6 + i * 2.5)
    add_card(s, cx, Inches(2.6), Inches(2.35), Inches(0.9), accent=OR)
    add_text(s, cx + Inches(0.25), Inches(2.7), Inches(2.0), Inches(0.3),
             name, size=11, bold=True, color=OR)
    add_text(s, cx + Inches(0.25), Inches(3.05), Inches(2.0), Inches(0.3),
             desc, size=10, color=TX2)

# 미션 화면 두 칸
add_text(s, Inches(0.6), Inches(3.8), Inches(12), Inches(0.4),
         "미션 슬라이드 구성 (두 칸)", size=14, bold=True, color=BLM)
# 왼쪽
add_card(s, Inches(0.6), Inches(4.3), Inches(5.0), Inches(2.8), accent=OR)
add_text(s, Inches(0.85), Inches(4.4), Inches(4.5), Inches(0.3),
         "왼쪽 35% — 안내 패널", size=12, bold=True, color=OR)
left_items = ["GOAL — 이 미션의 목적", "INPUT — 내가 입력할 것",
              "OUTPUT — 나와야 할 산출물", "프롬프트 + 힌트 + 모범답안",
              "필수/도전 체크 + AI 채점"]
box = s.shapes.add_textbox(Inches(0.85), Inches(4.8), Inches(4.5), Inches(2.2))
tf = box.text_frame; tf.word_wrap = True; tf.margin_left = tf.margin_right = Emu(0)
for i, t in enumerate(left_items):
    bullet(tf, t, size=10, color=TX, first=(i == 0))

# 오른쪽
add_card(s, Inches(5.8), Inches(4.3), Inches(7.1), Inches(2.8), accent=GD)
add_text(s, Inches(6.05), Inches(4.4), Inches(6.5), Inches(0.3),
         "오른쪽 65% — 터미널", size=12, bold=True, color=GD)
right_items = ["Claude Code가 이미 설치·준비됨",
               "우측 하단 AI 조수가 진행 코칭",
               "파일 생성 시 토스트 알림 + 미션 클리어"]
box = s.shapes.add_textbox(Inches(6.05), Inches(4.8), Inches(6.5), Inches(2.0))
tf = box.text_frame; tf.word_wrap = True; tf.margin_left = tf.margin_right = Emu(0)
for i, t in enumerate(right_items):
    bullet(tf, t, size=10, color=TX, first=(i == 0))

# ═══════════════════════════════════════════════════════
# Slide 8 — 오늘 배울 5가지 핵심 도구
# ═══════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
add_bg(s)
add_header(s, "TOOLS", "오늘 배울 5가지 핵심 도구")
add_text(s, Inches(0.6), Inches(2.1), Inches(12), Inches(0.5),
         "이 5가지를 순서대로 만들면 업무 자동화가 완성됩니다. 체험(Part 3)에서 하나씩, 실습(Part 4)에서 본인 업무에 적용.",
         size=14, color=TX2)

tools5 = [
    ("1", "🗺️", "Plan 모드", "먼저 설계 → 검토 → 승인\n엉뚱한 결과 방지", BL),
    ("2", "📋", "CLAUDE.md", "프로젝트 규칙서\n한 번 쓰면 매번 자동 적용", OR),
    ("3", "🎯", "Skill", "업무 표준 절차서\n같은 품질 결과 반복 가능", GD_D),
    ("4", "⌨️", "Command", "단축 명령어\n/ai-plan 한 줄이면 전체 실행", BLM),
    ("5", "🛡️", "Hook", "자동 안전장치\n개인정보 차단 등 예외 없이", WN),
]
for i, (num, icon, name, desc, accent) in enumerate(tools5):
    cy = Inches(2.75 + i * 0.92)
    add_card(s, Inches(0.6), cy, Inches(12), Inches(0.82), accent=accent)
    # number circle
    circle = s.shapes.add_shape(MSO_SHAPE.OVAL,
                                  Inches(0.85), cy + Inches(0.15),
                                  Inches(0.5), Inches(0.5))
    circle.fill.solid()
    circle.fill.fore_color.rgb = accent
    circle.line.fill.background()
    add_text(s, Inches(0.85), cy + Inches(0.15), Inches(0.5), Inches(0.5),
             num, size=16, bold=True, color=WHITE,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, Inches(1.55), cy + Inches(0.08), Inches(0.5), Inches(0.6),
             icon, size=22, color=accent, anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, Inches(2.15), cy + Inches(0.12), Inches(2.5), Inches(0.4),
             name, size=15, bold=True, color=accent)
    add_text(s, Inches(5.0), cy + Inches(0.12), Inches(7.5), Inches(0.6),
             desc.replace("\n", " — "), size=12, color=TX)

# ═══════════════════════════════════════════════════════
# Slide 9 — 2.1 시작하기
# ═══════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
add_bg(s)
add_header(s, "2.1", "Claude Code 시작하기")

add_text(s, Inches(0.6), Inches(2.2), Inches(12), Inches(0.5),
         "Claude Code = 나만의 AI 비서. 터미널 안에서 자연어로 말하면 행동합니다.",
         size=15, color=TX2)

# 큰 코드 블록
add_card(s, Inches(0.6), Inches(3.0), Inches(12), Inches(1.2), accent=OR)
add_text(s, Inches(1.0), Inches(3.35), Inches(11), Inches(0.5),
         "claude", size=28, bold=True, color=OR, font=MONO)

add_text(s, Inches(0.6), Inches(4.5), Inches(12), Inches(0.4),
         "이게 끝입니다. Claude Code가 시작되고 입력 프롬프트(>)가 나타납니다.",
         size=14, color=GD)

# 5가지 할 수 있는 일
items = [
    "폴더를 살펴보고 파일을 읽습니다",
    "새 파일을 만들거나 기존 파일을 수정합니다",
    "파이썬·노드 같은 코드를 직접 실행합니다",
    "인터넷을 검색해 최신 정보를 가져옵니다",
    "위 과정을 자율적으로 여러 단계 이어 실행합니다",
]
box = s.shapes.add_textbox(Inches(0.9), Inches(5.2), Inches(11.5), Inches(2.0))
tf = box.text_frame; tf.word_wrap = True; tf.margin_left = tf.margin_right = Emu(0)
for i, item in enumerate(items):
    bullet(tf, f"{i+1}. {item}", size=12, color=TX, first=(i == 0))

# ═══════════════════════════════════════════════════════
# Slide 10 — 좋은 프롬프트 쓰는 법 (6원칙 + Before/After)
# ═══════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
add_bg(s)
add_header(s, "PROMPT", "좋은 프롬프트 쓰는 법")
add_text(s, Inches(0.6), Inches(2.1), Inches(12), Inches(0.4),
         "이 6가지를 챙기면 AI가 한 번에 원하는 결과를 줍니다",
         size=14, color=TX2)

# 6원칙 (3x2 grid)
principles = [
    ("1", "주제 + 맥락", "무엇을, 왜", OR),
    ("2", "형식 + 구조", "파일 형태, 섹션", BLM),
    ("3", "단계별 분리", "한 번에 다 말고 나눠서", GD_D),
    ("4", "예시 제공", "원하는 형태 보여주기", WN),
    ("5", "역할 부여", "AI 에게 관점 지정", BL),
    ("6", "제한 사항", "하지 말아야 할 것", BAD),
]
for i, (num, name, desc, accent) in enumerate(principles):
    row, col = divmod(i, 3)
    cx = Inches(0.6 + col * 4.15)
    cy = Inches(2.6 + row * 1.15)
    add_card(s, cx, cy, Inches(3.95), Inches(0.95), accent=accent)
    circle = s.shapes.add_shape(MSO_SHAPE.OVAL, cx + Inches(0.2), cy + Inches(0.2), Inches(0.45), Inches(0.45))
    circle.fill.solid(); circle.fill.fore_color.rgb = accent; circle.line.fill.background()
    add_text(s, cx + Inches(0.2), cy + Inches(0.2), Inches(0.45), Inches(0.45),
             num, size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, cx + Inches(0.8), cy + Inches(0.15), Inches(3.0), Inches(0.4),
             name, size=14, bold=True, color=accent)
    add_text(s, cx + Inches(0.8), cy + Inches(0.5), Inches(3.0), Inches(0.4),
             desc, size=11, color=TX3)

# Before → After
add_card(s, Inches(0.6), Inches(5.2), Inches(5.5), Inches(1.8), accent=BAD)
add_text(s, Inches(0.85), Inches(5.3), Inches(5), Inches(0.35),
         "BEFORE", size=11, bold=True, color=BAD)
add_text(s, Inches(0.85), Inches(5.7), Inches(5), Inches(0.5),
         '"보고서 써줘"', size=16, color=TX, font=MONO)
add_text(s, Inches(0.85), Inches(6.3), Inches(5), Inches(0.5),
         "6가지 전부 빠짐 → AI 가 추측", size=11, color=TX3)

add_text(s, Inches(6.2), Inches(5.8), Inches(0.5), Inches(0.5),
         "→", size=24, bold=True, color=OR, anchor=MSO_ANCHOR.MIDDLE)

add_card(s, Inches(6.85), Inches(5.2), Inches(5.9), Inches(1.8), accent=GD)
add_text(s, Inches(7.1), Inches(5.3), Inches(5.5), Inches(0.35),
         "AFTER", size=11, bold=True, color=GD)
add_text(s, Inches(7.1), Inches(5.7), Inches(5.5), Inches(0.9),
         "2025 퇴직연금 시장 현황을\n금감원 데이터로 조사. 6섹션 Word.\n"
         "outputs/ 저장. 3페이지.",
         size=11, color=TX)
add_text(s, Inches(7.1), Inches(6.5), Inches(5.5), Inches(0.4),
         "6원칙 충족 → 한 번에 OK", size=11, color=GD)

# ═══════════════════════════════════════════════════════
# Slide 11 — 2.3 Plan 모드
# ═══════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
add_bg(s)
add_header(s, "2.3", "Plan 모드 — 본격 작업 전에 설계부터")

add_card(s, Inches(0.6), Inches(2.2), Inches(12), Inches(1.4), accent=OR)
add_text(s, Inches(1.0), Inches(2.5), Inches(11), Inches(1.0),
         '/plan AI 추진 TF의 "AI 추진 계획 보고서"를 자동으로\n'
         '만드는 프로그램을 설계할 거야. 단계별 계획을 세워줘.',
         size=14, color=OR, font=MONO)

add_text(s, Inches(0.6), Inches(3.9), Inches(12), Inches(0.5),
         "/plan으로 시작하면 Claude는 즉시 실행하지 않고 4~5단계 계획을 먼저 제시합니다.",
         size=13, color=TX2)

# 효과 3가지
effects = [
    ("방향 조기 차단", "의도와 다른 방향이면\n실행 전에 막을 수 있습니다", OR),
    ("비용 예측", "실행에 얼마나 걸릴지,\n어떤 도구가 호출될지 미리 확인", BLM),
    ("사용자가 흐름 통제", "AI 자율 실행이 아니라\n사람이 검토 후 승인", GD),
]
for i, (name, desc, accent) in enumerate(effects):
    cx = Inches(0.6 + i * 4.15)
    add_card(s, cx, Inches(4.7), Inches(3.95), Inches(2.2), accent=accent)
    add_text(s, cx + Inches(0.25), Inches(4.85), Inches(3.5), Inches(0.4),
             name, size=14, bold=True, color=accent)
    add_text(s, cx + Inches(0.25), Inches(5.35), Inches(3.5), Inches(1.2),
             desc, size=12, color=TX)

# ═══════════════════════════════════════════════════════
# Slide 12 — 2.4 권한
# ═══════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
add_bg(s)
add_header(s, "2.4", "권한 시스템 — 매번 Yes를 안 누르는 법")

add_card(s, Inches(0.6), Inches(2.2), Inches(7.5), Inches(3.5), accent=OR)
add_text(s, Inches(0.85), Inches(2.35), Inches(7.0), Inches(0.4),
         ".claude/settings.local.json", size=14, bold=True, color=OR, font=MONO)
code = ('{\n'
        '  "permissions": {\n'
        '    "allow": [\n'
        '      "Bash(*)",\n'
        '      "Read(*)",\n'
        '      "Write(*)",\n'
        '      "Edit(*)"\n'
        '    ]\n'
        '  }\n'
        '}')
add_text(s, Inches(0.85), Inches(2.85), Inches(7.0), Inches(2.8),
         code, size=12, color=GD, font=MONO)

add_card(s, Inches(8.4), Inches(2.2), Inches(4.5), Inches(3.5), accent=WN)
add_text(s, Inches(8.65), Inches(2.35), Inches(4.0), Inches(0.4),
         "워크숍에서는", size=14, bold=True, color=GD)
add_text(s, Inches(8.65), Inches(2.85), Inches(4.0), Inches(0.5),
         "격리된 환경이라 안전합니다.\n마음 편히 Yes를 누르세요.",
         size=12, color=TX)
add_text(s, Inches(8.65), Inches(3.7), Inches(4.0), Inches(0.4),
         "실제 업무에서는", size=14, bold=True, color=WN)
add_text(s, Inches(8.65), Inches(4.2), Inches(4.0), Inches(1.2),
         "Bash(*) 대신 Bash(python *)\n특정 폴더로만 쓰기 제한\nrm -rf 등 절대 자동 허용 금지",
         size=12, color=TX2)

# ═══════════════════════════════════════════════════════
# Slide 13 — 2.5 CLAUDE.md
# ═══════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
add_bg(s)
add_header(s, "2.5", "CLAUDE.md — 프로젝트 규칙서")

add_card(s, Inches(0.6), Inches(2.2), Inches(6.5), Inches(3.5), accent=GD)
add_text(s, Inches(0.85), Inches(2.35), Inches(6.0), Inches(0.4),
         "구조 예시", size=14, bold=True, color=GD)
md_example = ("# 프로젝트\n"
              "- 팀: AI 추진 TF\n"
              "- 목적: AI 추진 계획 보고서 자동 생성\n\n"
              "## 코딩 규칙\n"
              "- 모든 응답, 주석, 문서는 한국어\n"
              "- 강조 색상: #F58220\n"
              "- 모든 보고서는 outputs/ 폴더에 저장\n\n"
              "## 주요 작업\n"
              "- 리서치 → 본문 작성 → 검토 → .docx 출력")
add_text(s, Inches(0.85), Inches(2.85), Inches(6.0), Inches(2.7),
         md_example, size=11, color=TX, font=MONO)

# 전후 비교
add_card(s, Inches(7.4), Inches(2.2), Inches(5.5), Inches(1.6), accent=BAD)
add_text(s, Inches(7.65), Inches(2.35), Inches(5.0), Inches(0.3),
         "CLAUDE.md가 없을 때", size=12, bold=True, color=BAD)
add_text(s, Inches(7.65), Inches(2.75), Inches(5.0), Inches(0.9),
         "영어로 답하거나, 저장 위치 모르거나,\n일반적 톤으로 작성",
         size=11, color=TX2)

add_card(s, Inches(7.4), Inches(4.1), Inches(5.5), Inches(1.6), accent=GD)
add_text(s, Inches(7.65), Inches(4.25), Inches(5.0), Inches(0.3),
         "CLAUDE.md가 있을 때", size=12, bold=True, color=GD)
add_text(s, Inches(7.65), Inches(4.65), Inches(5.0), Inches(0.9),
         "한국어 + outputs/ 폴더 +\n우리 팀 톤 + 브랜드 색상 적용",
         size=11, color=TX)

add_text(s, Inches(0.6), Inches(6.0), Inches(12), Inches(0.4),
         "CLAUDE.md = \"AI를 우리 팀에 입사시키는 신입사원 교육 매뉴얼\"",
         size=13, color=GD)

# ═══════════════════════════════════════════════════════
# Slide 14 — 2.6 Skill
# ═══════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
add_bg(s)
add_header(s, "2.6", "Skill — 재사용 가능한 업무 매뉴얼")

add_card(s, Inches(0.6), Inches(2.2), Inches(6.5), Inches(4.0), accent=WN)
add_text(s, Inches(0.85), Inches(2.35), Inches(6.0), Inches(0.4),
         ".claude/skills/ai-plan-report/SKILL.md", size=12, bold=True, color=WN, font=MONO)
skill_ex = ("---\n"
            "name: ai-plan-report\n"
            "description: AI 추진 계획 보고서 작성\n"
            "---\n\n"
            "## 입력\n- 부서명, 프로젝트 목표\n\n"
            "## 절차\n"
            "1. 추진 배경 정리 (3~5문장)\n"
            "2. 적용 사례 리서치 (3개 이상)\n"
            "3. 도입 로드맵 (단계별 일정)\n"
            "4. 위험·기대효과 (각 3개)\n"
            "5. 결론 + .docx 출력\n\n"
            "## 출력\n- outputs/ai-plan-{부서}-{날짜}.docx")
add_text(s, Inches(0.85), Inches(2.85), Inches(6.0), Inches(3.2),
         skill_ex, size=10, color=TX, font=MONO)

# 효과
effects_skill = [
    ("한 줄 호출", "\"AI 추진 보고서 만들어 줘\"\n한 줄이면 5단계 자동 실행", GD),
    ("품질 일관성", "누가 시키든 동일한 구조의\n결과물이 나옵니다", BLM),
    ("팀 공유", "Skill 파일을 깃에 올리면\n팀 전체가 같은 절차로", OR),
]
for i, (name, desc, accent) in enumerate(effects_skill):
    cy = Inches(2.2 + i * 1.35)
    add_card(s, Inches(7.4), cy, Inches(5.5), Inches(1.2), accent=accent)
    add_text(s, Inches(7.65), cy + Inches(0.15), Inches(5.0), Inches(0.3),
             name, size=13, bold=True, color=accent)
    add_text(s, Inches(7.65), cy + Inches(0.5), Inches(5.0), Inches(0.6),
             desc, size=11, color=TX)

# ═══════════════════════════════════════════════════════
# Slide 15 — 2.7 Command
# ═══════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
add_bg(s)
add_header(s, "2.7", "Command — 단축 명령어")

add_card(s, Inches(0.6), Inches(2.2), Inches(6.5), Inches(3.0), accent=BL)
add_text(s, Inches(0.85), Inches(2.35), Inches(6.0), Inches(0.4),
         ".claude/commands/ai-plan.md", size=12, bold=True, color=BL, font=MONO)
cmd_ex = ("---\n"
          'description: 주제를 입력하면 보고서 자동 실행\n'
          "argument-hint: [주제]\n"
          "allowed-tools: Read, Write, Bash, WebFetch\n"
          "---\n\n"
          '다음 단계를 순서대로 실행:\n'
          '1. ai-plan-report 스킬로 보고서 작성\n'
          "2. PowerPoint 생성\n"
          "3. compliance-check 스킬로 규제 검토\n"
          "4. outputs/ 폴더에 저장")
add_text(s, Inches(0.85), Inches(2.85), Inches(6.0), Inches(2.2),
         cmd_ex, size=10, color=TX, font=MONO)

add_card(s, Inches(7.4), Inches(2.2), Inches(5.5), Inches(1.4), accent=OR)
add_text(s, Inches(7.65), Inches(2.35), Inches(5.0), Inches(0.3),
         "사용 예", size=14, bold=True, color=OR)
add_text(s, Inches(7.65), Inches(2.85), Inches(5.0), Inches(0.5),
         "/report 2025 퇴직연금 시장 동향", size=14, color=OR, font=MONO)

effects_cmd = [
    ("반복 작업 단축", "자주 쓰는 워크플로우를 한 단어로"),
    ("공유 가능", "깃에 올리면 팀이 같은 명령어"),
    ("인자 지원", "$ARGUMENTS로 동적 입력"),
]
for i, (name, desc) in enumerate(effects_cmd):
    cy = Inches(3.9 + i * 0.85)
    add_card(s, Inches(7.4), cy, Inches(5.5), Inches(0.7), accent=GD)
    add_text(s, Inches(7.65), cy + Inches(0.1), Inches(2.0), Inches(0.3),
             name, size=11, bold=True, color=GD)
    add_text(s, Inches(9.7), cy + Inches(0.1), Inches(3.0), Inches(0.5),
             desc, size=11, color=TX2)

# ═══════════════════════════════════════════════════════
# Slide 16 — 2.8 Hook
# ═══════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
add_bg(s)
add_header(s, "2.8", "Hook — 자동 안전장치")

add_card(s, Inches(0.6), Inches(2.2), Inches(7.0), Inches(3.5), accent=BAD)
add_text(s, Inches(0.85), Inches(2.35), Inches(6.5), Inches(0.4),
         "PII 차단 Hook 설정 (settings.local.json)", size=12, bold=True, color=BAD)
hook_json = ('{\n'
             '  "hooks": {\n'
             '    "PreToolUse": [{\n'
             '      "matcher": "Write",\n'
             '      "hooks": [{\n'
             '        "type": "command",\n'
             '        "command": "bash .claude/hooks/check-pii.sh"\n'
             '      }]\n'
             '    }]\n'
             '  }\n'
             '}')
add_text(s, Inches(0.85), Inches(2.85), Inches(6.5), Inches(2.7),
         hook_json, size=11, color=TX, font=MONO)

# 효과
effects_hook = [
    ("컴플라이언스 안심", "Claude의 판단에 의존하지 않고\n기계적으로 차단합니다", GD),
    ("다양한 활용", "코드 포맷팅, 테스트 자동 실행,\n변경 로그 기록 등", BLM),
    ("타이밍 선택", "PreToolUse / PostToolUse /\nStop 등 선택 가능", OR),
]
for i, (name, desc, accent) in enumerate(effects_hook):
    cy = Inches(2.2 + i * 1.2)
    add_card(s, Inches(7.9), cy, Inches(5.0), Inches(1.05), accent=accent)
    add_text(s, Inches(8.15), cy + Inches(0.12), Inches(4.5), Inches(0.3),
             name, size=12, bold=True, color=accent)
    add_text(s, Inches(8.15), cy + Inches(0.45), Inches(4.5), Inches(0.5),
             desc, size=10, color=TX)

add_text(s, Inches(0.6), Inches(6.1), Inches(12), Inches(0.5),
         "핵심: 사용자가 명시한 안전장치는 절대 우회되지 않습니다.\n"
         "Claude가 아무리 똑똑해도, Hook이 막으면 실행 불가.",
         size=12, color=WN)

# ═══════════════════════════════════════════════════════
# Slide 17 — 2.9~2.12 컨텍스트 + 성공 패턴 + 바이브코딩
# ═══════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
add_bg(s)
add_header(s, "2.9-2.12", "컨텍스트 관리 + 성공 패턴 + 바이브코딩 주의")

# 컨텍스트 관리
add_card(s, Inches(0.6), Inches(2.2), Inches(4.0), Inches(2.5), accent=BLM)
add_text(s, Inches(0.85), Inches(2.35), Inches(3.5), Inches(0.3),
         "컨텍스트 관리", size=14, bold=True, color=BLM)
ctx_items = ["/clear — 대화 완전 초기화", "/compact — 요약해서 토큰 절약",
             "/memory — 영구 메모리 관리", "/context — 사용량 확인"]
box = s.shapes.add_textbox(Inches(0.85), Inches(2.8), Inches(3.5), Inches(1.8))
tf = box.text_frame; tf.word_wrap = True; tf.margin_left = tf.margin_right = Emu(0)
for i, item in enumerate(ctx_items):
    bullet(tf, item, size=11, color=TX, mono=True, first=(i == 0))

# 성공 패턴
add_card(s, Inches(4.8), Inches(2.2), Inches(4.0), Inches(2.5), accent=GD)
add_text(s, Inches(5.05), Inches(2.35), Inches(3.5), Inches(0.3),
         "성공 패턴", size=14, bold=True, color=GD)
success = ["1. 일단 만들어 봄 (60점 OK)",
           "2. 결과 보고 수정 요청",
           "3. 다시 만듦 (75점)",
           "4. 만족할 때까지 반복"]
box = s.shapes.add_textbox(Inches(5.05), Inches(2.8), Inches(3.5), Inches(1.8))
tf = box.text_frame; tf.word_wrap = True; tf.margin_left = tf.margin_right = Emu(0)
for i, item in enumerate(success):
    bullet(tf, item, size=11, color=TX, first=(i == 0))

# 도구 병렬 실행
add_card(s, Inches(9.0), Inches(2.2), Inches(3.9), Inches(2.5), accent=OR)
add_text(s, Inches(9.25), Inches(2.35), Inches(3.5), Inches(0.3),
         "도구 병렬 실행", size=14, bold=True, color=OR)
parallel = ["Read, Write, Edit, Bash, WebFetch",
            "동시 호출로 시간 절반",
            "의존성 있으면 자동 순차 실행"]
box = s.shapes.add_textbox(Inches(9.25), Inches(2.8), Inches(3.5), Inches(1.8))
tf = box.text_frame; tf.word_wrap = True; tf.margin_left = tf.margin_right = Emu(0)
for i, item in enumerate(parallel):
    bullet(tf, item, size=11, color=TX, first=(i == 0))

# 바이브코딩 주의사항 표
add_text(s, Inches(0.6), Inches(5.0), Inches(12), Inches(0.3),
         "바이브코딩 주의사항", size=14, bold=True, color=WN)
warn_h = ["위험", "방지법"]
warns = [
    ["할루시네이션", "출처 요구, 결과 직접 검증, 원본 확인"],
    ["보안 위험", "Hook으로 개인정보·키 차단, 권한 좁게"],
    ["저작권", "출처 명시, 허가된 자료만 사용"],
    ["잘못된 자동화", "결과 검토 후 사용, 중간 체크"],
    ["비밀 유출", "API 키 하드코딩 금지, 환경 변수"],
]
warn_cw = [3.0, 9.0]
tbl = s.shapes.add_table(
    len(warns) + 1, 2,
    Inches(0.6), Inches(5.4),
    Inches(sum(warn_cw)), Inches(0.32 * (len(warns) + 1)),
).table
for ci, h in enumerate(warn_h):
    tbl.columns[ci].width = Inches(warn_cw[ci])
    cell = tbl.cell(0, ci)
    cell.text = h
    for p in cell.text_frame.paragraphs:
        p.font.name = KO_FONT; p.font.size = Pt(10); p.font.bold = True; p.font.color.rgb = WHITE
    cell.fill.solid(); cell.fill.fore_color.rgb = NV
for ri, row in enumerate(warns):
    for ci, val in enumerate(row):
        cell = tbl.cell(ri + 1, ci)
        cell.text = val
        for p in cell.text_frame.paragraphs:
            p.font.name = KO_FONT; p.font.size = Pt(9); p.font.color.rgb = TX
        cell.fill.solid()
        cell.fill.fore_color.rgb = NV_D if ri % 2 == 0 else BG

# ═══════════════════════════════════════════════════════
# Slide 18 — Part 3 소개
# ═══════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
add_bg(s)
add_header(s, "PART 3", "기능 체험 — 8개의 미션")

add_text(s, Inches(0.6), Inches(2.1), Inches(12), Inches(0.5),
         '공통 주제: AI 추진 TF의 "AI 추진 계획 보고서 자동화"',
         size=15, bold=True, color=OR)

# 미션 구조 카드
add_text(s, Inches(0.6), Inches(2.8), Inches(12), Inches(0.3),
         "각 미션의 구조", size=14, bold=True, color=BLM)
struct = [
    ("GOAL", GD, "이 미션에서 배우는 것"),
    ("INPUT", OR, "무엇을 입력하는지"),
    ("OUTPUT", BLM, "어떤 파일이 나와야 하는지"),
]
for i, (name, accent, desc) in enumerate(struct):
    cx = Inches(0.6 + i * 4.15)
    add_card(s, cx, Inches(3.2), Inches(3.95), Inches(0.8), accent=accent)
    add_text(s, cx + Inches(0.3), Inches(3.3), Inches(1.5), Inches(0.3),
             name, size=12, bold=True, color=accent)
    add_text(s, cx + Inches(1.8), Inches(3.3), Inches(2.0), Inches(0.6),
             desc, size=11, color=TX)

# 추가 요소
add_card(s, Inches(0.6), Inches(4.3), Inches(12), Inches(2.8), accent=GD)
add_text(s, Inches(0.85), Inches(4.45), Inches(11.5), Inches(0.3),
         "자동 채점 + 체크리스트", size=14, bold=True, color=GD)
check_items = [
    "프롬프트 — 그대로 복사하거나 직접 작성",
    "✅ 필수 체크리스트 — 통과해야 미션 클리어",
    "⭐ 도전 체크리스트 — 선택, 더 좋은 결과물을 위한 추가 기준",
    "🤖 AI 채점 — 산출물을 AI가 읽고 항목별 점수 + 코멘트",
    "🤖 AI 조수 — 터미널 상단에서 실시간 진행 코칭",
    "📄 산출물 프리뷰 — 파일 생성 후 클릭하면 내용 직접 보기 + ⬇ 다운로드",
    "파일 생성 → 토스트 알림 → 체크 자동 통과 → 미션 클리어",
]
box = s.shapes.add_textbox(Inches(0.85), Inches(4.9), Inches(11.5), Inches(2.0))
tf = box.text_frame; tf.word_wrap = True; tf.margin_left = tf.margin_right = Emu(0)
for i, item in enumerate(check_items):
    bullet(tf, item, size=12, color=TX, first=(i == 0))

# ═══════════════════════════════════════════════════════
# Slide 19 — 미션 1-4
# ═══════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
add_bg(s)
add_header(s, "MISSION 1-4", "미션 1~4: Plan / 권한 / CLAUDE.md / Skill")

missions_1_4 = [
    ("M1", "Plan 모드", OR,
     '/plan AI 추진 TF의 "AI 추진 계획 보고서"를\n자동으로 만드는 프로그램을 설계할 거야.\n단계별 계획을 세우고 PLAN.md로 저장해줘.',
     "PLAN.md"),
    ("M2", "권한 설정", BLM,
     ".claude/settings.local.json 파일을 만들어서\nBash, Read, Write, Edit 도구를 자동 허용.",
     ".claude/settings.local.json"),
    ("M3", "CLAUDE.md", GD,
     '현재 폴더에 새 CLAUDE.md를 만들어줘.\n팀: AI 추진 TF, 프로젝트: AI 추진 계획.\n한국어, #F58220, 250단어 이내.',
     "CLAUDE.md"),
    ("M4", "Skill", WN,
     ".claude/skills/ai-plan-report/SKILL.md\n입력, 절차(5단계), 출력 형식 명시.\n250단어 이내.",
     ".claude/skills/.../SKILL.md"),
]
for i, (num, title, accent, prompt, target) in enumerate(missions_1_4):
    row, col = divmod(i, 2)
    cx = Inches(0.6 + col * 6.3)
    cy = Inches(2.2 + row * 2.55)
    add_card(s, cx, cy, Inches(6.1), Inches(2.35), accent=accent)
    add_chip(s, cx + Inches(0.2), cy + Inches(0.15), num, fill=accent, size=10)
    add_text(s, cx + Inches(1.0), cy + Inches(0.12), Inches(4.8), Inches(0.35),
             title, size=14, bold=True, color=accent)
    add_text(s, cx + Inches(0.25), cy + Inches(0.55), Inches(5.6), Inches(1.2),
             prompt, size=9, color=TX, font=MONO)
    add_text(s, cx + Inches(0.25), cy + Inches(1.85), Inches(5.6), Inches(0.3),
             f"목표: {target}", size=10, bold=True, color=TX2)

# ═══════════════════════════════════════════════════════
# Slide 20 — 미션 5-8
# ═══════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
add_bg(s)
add_header(s, "MISSION 5-8", "미션 5~8: Command / Hook / 컨텍스트 / 노트")

missions_5_8 = [
    ("M5", "Command", BL,
     '"/ai-plan" 슬래시 커맨드를\n.claude/commands/ai-plan.md로 만들어줘.\n인자로 부서명을 받아 ai-plan-report 호출.',
     ".claude/commands/ai-plan.md"),
    ("M6", "Hook", BAD,
     '보고서 저장 전에 개인정보\n(주민번호, 전화번호, 이메일) 차단 훅.\n.claude/hooks/check-pii.sh + Hook 등록.',
     ".claude/hooks/check-pii.sh"),
    ("M7", "컨텍스트 관리", BLM,
     '지금까지 AI 추진 계획 보고서 프로젝트로\n진행한 사항들을 기억해줘.\n그 다음 /compact 실행해줘.',
     "대화 요약 + 컨텍스트 압축"),
    ("M8", "노트 테이킹", GD_D,
     '이 프로젝트는 AI 추진 TF의\n"AI 추진 계획 보고서 자동화"야.\n이 정보를 메모리에 기억해줘.',
     "~/.claude.json 영구 메모리"),
]
for i, (num, title, accent, prompt, target) in enumerate(missions_5_8):
    row, col = divmod(i, 2)
    cx = Inches(0.6 + col * 6.3)
    cy = Inches(2.2 + row * 2.55)
    add_card(s, cx, cy, Inches(6.1), Inches(2.35), accent=accent)
    add_chip(s, cx + Inches(0.2), cy + Inches(0.15), num, fill=accent, size=10)
    add_text(s, cx + Inches(1.0), cy + Inches(0.12), Inches(4.8), Inches(0.35),
             title, size=14, bold=True, color=accent)
    add_text(s, cx + Inches(0.25), cy + Inches(0.55), Inches(5.6), Inches(1.2),
             prompt, size=9, color=TX, font=MONO)
    add_text(s, cx + Inches(0.25), cy + Inches(1.85), Inches(5.6), Inches(0.3),
             f"목표: {target}", size=10, bold=True, color=TX2)

# ═══════════════════════════════════════════════════════
# Slide 21 — Part 4 소개
# ═══════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
add_bg(s)
add_header(s, "PART 4", "실습 프로젝트 — 7단계 워크플로우")

add_text(s, Inches(0.6), Inches(2.1), Inches(12), Inches(0.5),
         "모든 도구를 한 번에 활용해서, 본인 업무에 맞는 자동화를 처음부터 끝까지 완성합니다.",
         size=14, color=TX2)

steps_part4 = [
    ("1", "Plan", "만들 프로그램의 단계별 설계\n(PLAN.md 저장)", OR),
    ("2", "CLAUDE.md", "본인 업무에 맞는\n프로젝트 규칙서 작성", GD),
    ("3", "Skill", "핵심 업무 절차를\n스킬로 정리", WN),
    ("4", "Command", "자주 쓸 워크플로우를\n단축 명령어로", BL),
    ("5", "Hook", "안전장치 추가", BAD),
    ("6", "실행", "자기 부서 주제로\n직접 돌려 보기", BLM),
    ("7", "웹페이지", "완성 기능을 HTML로\n다른 사람도 쓸 수 있게", GD_D),
]
for i, (num, name, desc, accent) in enumerate(steps_part4):
    if i < 4:
        cx = Inches(0.6 + i * 3.2)
        cy = Inches(2.8)
    else:
        cx = Inches(0.6 + (i - 4) * 4.15)
        cy = Inches(5.0)
    w = Inches(3.0) if i < 4 else Inches(3.95)
    h = Inches(1.9) if i < 4 else Inches(1.9)
    add_card(s, cx, cy, w, h, accent=accent)
    circle = s.shapes.add_shape(MSO_SHAPE.OVAL, cx + Inches(0.2), cy + Inches(0.25), Inches(0.5), Inches(0.5))
    circle.fill.solid(); circle.fill.fore_color.rgb = accent
    circle.line.fill.background()
    add_text(s, cx + Inches(0.2), cy + Inches(0.25), Inches(0.5), Inches(0.5),
             num, size=16, bold=True, color=WHITE, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, cx + Inches(0.85), cy + Inches(0.25), w - Inches(1.1), Inches(0.35),
             name, size=13, bold=True, color=accent)
    add_text(s, cx + Inches(0.25), cy + Inches(0.85), w - Inches(0.4), Inches(0.9),
             desc, size=10, color=TX)

# ═══════════════════════════════════════════════════════
# Slide 22 — 완성하면
# ═══════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
add_bg(s)
add_header(s, "4.2", "완성하면 무엇이 손에 남을까요")

results = [
    ("Plan 모드", "설계 → 검토 → 승인 단계", OR),
    ("CLAUDE.md", "매번 자동 적용되는 프로젝트 규칙", GD),
    ("Skill", "재사용 가능한 표준 업무 매뉴얼", WN),
    ("Command", "/명령어 한 줄로 전체 워크플로우 호출", BL),
    ("Hook", "예외 없는 자동 안전장치", BAD),
    ("Web UI", "다른 사람도 쓸 수 있는 HTML 인터페이스", BLM),
]
for i, (name, desc, accent) in enumerate(results):
    row, col = divmod(i, 3)
    cx = Inches(0.6 + col * 4.15)
    cy = Inches(2.3 + row * 2.4)
    add_card(s, cx, cy, Inches(3.95), Inches(2.1), accent=accent)
    add_text(s, cx + Inches(0.3), cy + Inches(0.2), Inches(3.5), Inches(0.4),
             name, size=16, bold=True, color=accent)
    add_text(s, cx + Inches(0.3), cy + Inches(0.7), Inches(3.5), Inches(1.0),
             desc, size=13, color=TX)

add_text(s, Inches(0.6), Inches(7.0), Inches(12), Inches(0.4),
         "이 여섯 가지 조합이 여러분의 첫 \"AI 자동화 도구\"입니다.",
         size=13, color=GD)

# ═══════════════════════════════════════════════════════
# Slide 23 — A.1 CLI 설치
# ═══════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
add_bg(s)
add_header(s, "A.1", "Claude Code (CLI) 개인 PC 설치")

# 사전 준비
cards_a1 = [
    ("Node.js 18 이상", "node --version", "nodejs.org 에서 LTS 다운로드"),
    ("Git", "git --version", "Windows: Git for Windows (Git Bash 포함)"),
    ("Anthropic 계정", "claude.ai 가입", "Pro / Max 구독 권장"),
]
y = Inches(2.2)
for i, (title, cmd, desc) in enumerate(cards_a1):
    cy = y + Inches(i * 1.1)
    add_card(s, Inches(0.6), cy, Inches(12), Inches(0.9), accent=OR)
    add_text(s, Inches(0.9), cy + Inches(0.12), Inches(3.5), Inches(0.35),
             title, size=14, bold=True, color=OR)
    add_text(s, Inches(0.9), cy + Inches(0.5), Inches(3.5), Inches(0.3),
             cmd, size=12, color=TX, font=MONO)
    add_text(s, Inches(5.0), cy + Inches(0.3), Inches(7.5), Inches(0.4),
             desc, size=12, color=TX2)

# 설치 명령
add_card(s, Inches(0.6), Inches(5.6), Inches(12), Inches(1.2), accent=GD)
add_text(s, Inches(0.85), Inches(5.7), Inches(11.5), Inches(0.3),
         "설치 (모든 OS 공통)", size=12, bold=True, color=GD)
add_text(s, Inches(0.85), Inches(6.1), Inches(11.5), Inches(0.5),
         "npm install -g @anthropic-ai/claude-code", size=18, bold=True, color=OR, font=MONO)

# ═══════════════════════════════════════════════════════
# Slide 24 — A.2 Desktop 설치
# ═══════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
add_bg(s)
add_header(s, "A.2", "Claude Desktop (GUI 앱) 설치")

add_text(s, Inches(0.6), Inches(2.2), Inches(12), Inches(0.5),
         "claude.ai 다운로드 페이지에서 본인 OS에 맞는 설치 파일을 받으세요.",
         size=14, color=TX2)

# OS 카드
for i, (os_name, fname, accent) in enumerate([
    ("Windows", "Claude-Setup.exe", BLM),
    ("macOS", "Claude.dmg", TX2),
    ("Linux", ".deb / AppImage", GD),
]):
    cx = Inches(0.6 + i * 4.15)
    add_card(s, cx, Inches(3.0), Inches(3.95), Inches(1.8), accent=accent)
    add_text(s, cx + Inches(0.3), Inches(3.15), Inches(3.5), Inches(0.5),
             os_name, size=20, bold=True, color=accent)
    add_text(s, cx + Inches(0.3), Inches(3.75), Inches(3.5), Inches(0.4),
             fname, size=13, color=TX, font=MONO)
    add_text(s, cx + Inches(0.3), Inches(4.2), Inches(3.5), Inches(0.4),
             "다음 > 다음 > 설치 완료", size=11, color=TX3)

# 첫 실행
add_card(s, Inches(0.6), Inches(5.2), Inches(12), Inches(1.6), accent=OR)
add_text(s, Inches(0.85), Inches(5.35), Inches(11.5), Inches(0.3),
         "첫 실행", size=14, bold=True, color=OR)
add_text(s, Inches(0.85), Inches(5.75), Inches(11.5), Inches(0.4),
         "앱 실행 → Claude 계정 로그인 → 바로 채팅 시작",
         size=13, color=TX)
add_text(s, Inches(0.85), Inches(6.2), Inches(11.5), Inches(0.4),
         "MCP 연결 (선택): 메뉴 → 설정 → MCP에서 Filesystem, GitHub, Slack 등 커넥터 연결 가능",
         size=11, color=TX3)

# ═══════════════════════════════════════════════════════
# Slide 25 — A.3 비교
# ═══════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
add_bg(s)
add_header(s, "A.3", "CLI vs Desktop 비교")

cmp_h = ["구분", "Claude Code (CLI)", "Claude Desktop (GUI)"]
cmp_rows = [
    ["인터페이스", "터미널 (검은 창)", "채팅 앱 (윈도우/맥 GUI)"],
    ["학습 곡선", "약간 있음", "매우 쉬움 (카톡 수준)"],
    ["파일 생성·수정", "자유롭게 가능", "제한적 (대화 첨부만)"],
    ["코드 직접 실행", "가능 (Bash, Python 등)", "불가"],
    ["CLAUDE.md / Skill / Hook", "모두 사용 가능", "사용 불가"],
    ["자동화 워크플로우", "핵심 기능", "거의 불가"],
    ["적합한 용도", "보고서 자동화, 분석, 코드", "일상 질문, 글쓰기, 아이디어"],
]
cmp_cw = [3.0, 4.5, 4.5]
tbl = s.shapes.add_table(
    len(cmp_rows) + 1, 3,
    Inches(0.6), Inches(2.2),
    Inches(sum(cmp_cw)), Inches(0.4 * (len(cmp_rows) + 1)),
).table
for ci, h in enumerate(cmp_h):
    tbl.columns[ci].width = Inches(cmp_cw[ci])
    cell = tbl.cell(0, ci)
    cell.text = h
    for p in cell.text_frame.paragraphs:
        p.font.name = KO_FONT; p.font.size = Pt(11); p.font.bold = True; p.font.color.rgb = WHITE
    cell.fill.solid(); cell.fill.fore_color.rgb = NV
for ri, row in enumerate(cmp_rows):
    for ci, val in enumerate(row):
        cell = tbl.cell(ri + 1, ci)
        cell.text = val
        for p in cell.text_frame.paragraphs:
            p.font.name = KO_FONT; p.font.size = Pt(10); p.font.color.rgb = TX
        cell.fill.solid()
        cell.fill.fore_color.rgb = NV_D if ri % 2 == 0 else BG

add_text(s, Inches(0.6), Inches(5.7), Inches(12), Inches(0.6),
         "가벼운 질문·글쓰기 → Desktop\n보고서 자동화·파일 작업 → Code\n"
         "처음에는 둘 다 설치해 두고 상황에 맞게 골라 쓰시는 것이 효율적입니다.",
         size=13, color=GD)

# ═══════════════════════════════════════════════════════
# Slide 26 — A.4 실행 흐름
# ═══════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
add_bg(s)
add_header(s, "A.4", "실행 흐름과 자주 쓰는 명령")

# 6단계 흐름
flow = [
    ("1", "작업 폴더 이동", "cd ~/projects/my-task"),
    ("2", "Claude Code 실행", "claude"),
    ("3", "자연어 지시", "> *.csv 합쳐서 summary.xlsx 만들어줘"),
    ("4", "Claude 작업", "→ 결과 파일 생성"),
    ("5", "수정 요청", "> 매출 컬럼에 천 단위 콤마 추가해줘"),
    ("6", "종료", "Ctrl+C 두 번 또는 /exit"),
]
for i, (num, name, desc) in enumerate(flow):
    row, col = divmod(i, 3)
    cx = Inches(0.6 + col * 4.15)
    cy = Inches(2.2 + row * 1.35)
    add_card(s, cx, cy, Inches(3.95), Inches(1.15), accent=OR)
    add_chip(s, cx + Inches(0.15), cy + Inches(0.15), num, fill=OR, size=10)
    add_text(s, cx + Inches(0.65), cy + Inches(0.12), Inches(3.1), Inches(0.3),
             name, size=11, bold=True, color=OR)
    add_text(s, cx + Inches(0.25), cy + Inches(0.5), Inches(3.5), Inches(0.5),
             desc, size=10, color=TX, font=MONO)

# 명령어 표
cmd_h = ["명령", "효과"]
cmds = [
    ["/login", "Claude 계정 로그인"],
    ["/clear", "현재 대화 초기화"],
    ["/compact", "대화 요약 (긴 대화 시)"],
    ["/memory", "영구 메모리 관리"],
    ["/context", "컨텍스트 사용량 확인"],
    ["/usage", "Claude Max 사용량 확인"],
    ["/help", "명령어 목록"],
    ["/exit", "종료"],
]
cmd_cw2 = [2.5, 9.5]
tbl = s.shapes.add_table(
    len(cmds) + 1, 2,
    Inches(0.6), Inches(5.1),
    Inches(sum(cmd_cw2)), Inches(0.28 * (len(cmds) + 1)),
).table
for ci, h in enumerate(cmd_h):
    tbl.columns[ci].width = Inches(cmd_cw2[ci])
    cell = tbl.cell(0, ci)
    cell.text = h
    for p in cell.text_frame.paragraphs:
        p.font.name = KO_FONT; p.font.size = Pt(9); p.font.bold = True; p.font.color.rgb = WHITE
    cell.fill.solid(); cell.fill.fore_color.rgb = NV
for ri, row in enumerate(cmds):
    for ci, val in enumerate(row):
        cell = tbl.cell(ri + 1, ci)
        cell.text = val
        for p in cell.text_frame.paragraphs:
            p.font.name = MONO if ci == 0 else KO_FONT
            p.font.size = Pt(9); p.font.color.rgb = TX
        cell.fill.solid()
        cell.fill.fore_color.rgb = NV_D if ri % 2 == 0 else BG

# ═══════════════════════════════════════════════════════
# Slide 27 — A.5 안전
# ═══════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
add_bg(s)
add_header(s, "A.5", "안전하게 쓰기 위한 주의사항")

safety = [
    ("권한과 폴더", WN, [
        "작업 폴더 분리: 회사/개인 자료를 폴더 단위로 분리",
        "권한 좁히기: Bash(*) 대신 Bash(python *) 부터",
        "CLAUDE.md에 보안 규칙 명시",
    ]),
    ("위험한 명령", BAD, [
        "rm -rf (영구 삭제) — 절대 자동 허용 금지",
        "git push --force (팀원 코드 덮어쓰기)",
        "sudo (관리자 권한) / curl, wget (외부 호출)",
    ]),
    ("데이터와 개인정보", BLM, [
        "회사 내부 자료 → 회사 정책 먼저 확인",
        "주민번호·전화번호·이메일 → Hook 자동 차단",
        "삭제·덮어쓰기 작업 → 백업 먼저",
    ]),
    ("비용 관리", GD, [
        "Pro/Max 구독 → 월 정액으로 충분",
        "API 키 직접 사용 → 토큰 단위 과금 주의",
    ]),
]
for i, (title, accent, items) in enumerate(safety):
    row, col = divmod(i, 2)
    cx = Inches(0.6 + col * 6.35)
    cy = Inches(2.2 + row * 2.5)
    h = Inches(2.25)
    add_card(s, cx, cy, Inches(6.1), h, accent=accent)
    add_text(s, cx + Inches(0.25), cy + Inches(0.15), Inches(5.5), Inches(0.35),
             title, size=14, bold=True, color=accent)
    box = s.shapes.add_textbox(cx + Inches(0.25), cy + Inches(0.55), Inches(5.6), h - Inches(0.7))
    tf = box.text_frame; tf.word_wrap = True; tf.margin_left = tf.margin_right = Emu(0)
    for j, item in enumerate(items):
        bullet(tf, item, size=11, color=TX, first=(j == 0))

# ═══════════════════════════════════════════════════════
# Slide 28 — 끝
# ═══════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
add_bg(s)
add_rect(s, 0, 0, Inches(5), SLIDE_H, fill=OR_D)
add_text(s, Inches(0.6), Inches(0.6), Inches(4), Inches(0.5),
         "MIRAE ASSET LIFE", size=14, bold=True, color=WHITE)
add_text(s, Inches(0.6), Inches(1.0), Inches(4), Inches(0.4),
         "Vibe Coding Workshop", size=12, color=WHITE)

add_text(s, Inches(5.6), Inches(2.5), Inches(7.4), Inches(1.5),
         "감사합니다", size=48, bold=True, color=WHITE)
add_text(s, Inches(5.6), Inches(4.2), Inches(7.4), Inches(0.6),
         "오늘 배운 한 가지를 일주일 안에\n본인 업무에 적용해보세요.",
         size=18, color=TX2)
add_text(s, Inches(5.6), Inches(5.4), Inches(7.4), Inches(0.5),
         "미래에셋생명  ·  Claude Code  ·  2026",
         size=12, color=TX3, font=MONO)

# ═══════════════════════════════════════════════════════
# 저장
# ═══════════════════════════════════════════════════════
prs.save(str(OUT))
print(f"OK -> {OUT} ({OUT.stat().st_size // 1024} KB, {len(prs.slides)} slides)")
