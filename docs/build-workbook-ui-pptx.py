#!/usr/bin/env python3
"""워크북 사용 가이드 (신규 UI) — PPTX.

handout-workbook-ui.md 기준으로 13장 슬라이드.
브랜드: 미래에셋생명 · 오렌지 #F58220 + 짙은 네이비.
build-home-practice-pptx.py 의 디자인 토큰/헬퍼와 동일.
"""
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

DOCS = Path(__file__).parent
OUT = DOCS / "workbook-ui.pptx"

# ── 브랜드 색상 ────────────────────────────────────────
OR = RGBColor(0xF5, 0x82, 0x20)
OR_D = RGBColor(0xCB, 0x60, 0x15)
NV = RGBColor(0x04, 0x3B, 0x72)
NV_D = RGBColor(0x06, 0x1E, 0x30)
BG = RGBColor(0x04, 0x18, 0x28)
TX = RGBColor(0xE5, 0xE8, 0xEC)
TX2 = RGBColor(0x8D, 0xA0, 0xB8)
TX3 = RGBColor(0x5A, 0x7A, 0x98)
GD = RGBColor(0x86, 0xEF, 0xAC)
GD_D = RGBColor(0x10, 0xB9, 0x81)
BAD = RGBColor(0xFC, 0xA5, 0xA5)
WN = RGBColor(0xFB, 0xBF, 0x24)
BL = RGBColor(0x60, 0xA5, 0xFA)
BLM = RGBColor(0x38, 0xBD, 0xF8)
BD = RGBColor(0x0A, 0x30, 0x50)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BLACK = RGBColor(0x00, 0x00, 0x00)

KO_FONT = "Pretendard"
MONO = "JetBrains Mono"

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
SLIDE_W = prs.slide_width
SLIDE_H = prs.slide_height
BLANK = prs.slide_layouts[6]


# ── 헬퍼 ──────────────────────────────────────────────
def add_bg(slide, color=BG):
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H)
    bg.fill.solid()
    bg.fill.fore_color.rgb = color
    bg.line.fill.background()
    return bg


def add_text(slide, x, y, w, h, text, *, size=18, color=TX, bold=False,
             font=KO_FONT, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP):
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Emu(0)
    tf.margin_top = tf.margin_bottom = Emu(0)
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return box


def add_rect(slide, x, y, w, h, *, fill=None, line=None, line_w=None, radius=False):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE,
        x, y, w, h,
    )
    if fill is not None:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    else:
        shape.fill.background()
    if line is not None:
        shape.line.color.rgb = line
        if line_w:
            shape.line.width = line_w
    else:
        shape.line.fill.background()
    return shape


def add_chip(slide, x, y, text, *, fill=OR, fg=WHITE, size=11, padx=8):
    w = Emu(int(len(text) * 110000 + padx * 18000))
    h = Inches(0.32)
    chip = add_rect(slide, x, y, w, h, fill=fill, radius=True)
    chip.line.fill.background()
    tf = chip.text_frame
    tf.margin_left = tf.margin_right = Emu(60000)
    tf.margin_top = tf.margin_bottom = Emu(0)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = text
    r.font.name = KO_FONT
    r.font.size = Pt(size)
    r.font.bold = True
    r.font.color.rgb = fg
    return chip


def add_header(slide, kicker, title):
    add_chip(slide, Inches(0.6), Inches(0.5), kicker)
    add_text(
        slide, Inches(0.6), Inches(0.95), Inches(12), Inches(0.9),
        title, size=30, bold=True, color=WHITE,
    )
    line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0.6), Inches(1.85), Inches(0.7), Emu(40000),
    )
    line.fill.solid()
    line.fill.fore_color.rgb = OR
    line.line.fill.background()
    add_text(
        slide, Inches(12.5), Inches(7.05), Inches(0.7), Inches(0.3),
        f"{len(prs.slides) + 1:02d}",
        size=10, color=TX3, font=MONO, align=PP_ALIGN.RIGHT,
    )


def add_card(slide, x, y, w, h, *, accent=OR):
    card = add_rect(slide, x, y, w, h, fill=NV_D, line=BD, line_w=Pt(0.75), radius=True)
    add_rect(slide, x, y, Inches(0.06), h, fill=accent)
    return card


def bullet(tf, text, *, size=13, color=TX, bold=False, indent=0, mono=False, first=False):
    p = tf.paragraphs[0] if first else tf.add_paragraph()
    p.alignment = PP_ALIGN.LEFT
    p.level = indent
    p.space_after = Pt(3)
    r = p.add_run()
    r.text = text
    r.font.name = MONO if mono else KO_FONT
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.color.rgb = color
    return p


def card_with_bullets(slide, x, y, w, h, accent, title, bullets, *, title_size=16):
    add_card(slide, x, y, w, h, accent=accent)
    add_text(slide, x + Inches(0.3), y + Inches(0.18), w - Inches(0.4), Inches(0.4),
             title, size=title_size, bold=True, color=accent)
    box = slide.shapes.add_textbox(
        x + Inches(0.3), y + Inches(0.65),
        w - Inches(0.5), h - Inches(0.8),
    )
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Emu(0)
    for i, ln in enumerate(bullets):
        bullet(tf, ln, size=12, color=TX, first=(i == 0))


# ═══════════════════════════════════════════════════════
# Slide 1 — 표지
# ═══════════════════════════════════════════════════════
s1 = prs.slides.add_slide(BLANK)
add_bg(s1)
add_rect(s1, 0, 0, Inches(5), SLIDE_H, fill=OR_D)
add_text(s1, Inches(0.6), Inches(0.6), Inches(4), Inches(0.5),
         "MIRAE ASSET LIFE", size=14, bold=True, color=WHITE)
add_text(s1, Inches(0.6), Inches(1.0), Inches(4), Inches(0.4),
         "Vibe Coding Workshop", size=12, color=WHITE)

add_text(s1, Inches(5.6), Inches(1.7), Inches(7.4), Inches(0.5),
         "WORKBOOK UI GUIDE", size=12, bold=True, color=OR)
add_text(s1, Inches(5.6), Inches(2.2), Inches(7.4), Inches(2.2),
         "워크북 사용\n가이드", size=44, bold=True, color=WHITE)
add_text(s1, Inches(5.6), Inches(5.0), Inches(7.4), Inches(0.6),
         "안내 보강 · AI 채점 · AI 조수 · 파일 가시화\n— 학습자 + 강사용",
         size=18, color=TX2)
add_text(s1, Inches(5.6), Inches(6.6), Inches(7.4), Inches(0.4),
         "미래에셋생명  ·  2026", size=11, color=TX3, font=MONO)

# ═══════════════════════════════════════════════════════
# Slide 2 — 화면 한눈에 (좌/우 도식)
# ═══════════════════════════════════════════════════════
s2 = prs.slides.add_slide(BLANK)
add_bg(s2)
add_header(s2, "OVERVIEW", "워크북 화면 한눈에")
add_text(s2, Inches(0.6), Inches(2.1), Inches(12), Inches(0.6),
         "미션 슬라이드는 두 칸 — 왼쪽 안내 + 오른쪽 터미널.\n오른쪽 하단엔 항상 AI 조수가 떠 있어 학습자 진행을 코칭합니다.",
         size=14, color=TX2)

# 왼쪽 패널 (briefing)
add_card(s2, Inches(0.6), Inches(3.5), Inches(4.4), Inches(3.4))
add_text(s2, Inches(0.85), Inches(3.65), Inches(4), Inches(0.4),
         "왼쪽 (35%)", size=11, bold=True, color=OR)
add_text(s2, Inches(0.85), Inches(3.95), Inches(4), Inches(0.4),
         "안내 / 체크리스트", size=15, bold=True, color=WHITE)
items_l = [
    "🎯 GOAL — 이 미션의 목적",
    "📥 INPUT — 내가 입력할 것",
    "📤 OUTPUT — 나와야 할 산출물",
    "💬 프롬프트 / 힌트 / 모범답안",
    "✅ 필수 체크 / ⭐ 도전 체크",
    "🤖 AI 채점 받기",
]
box = s2.shapes.add_textbox(Inches(0.85), Inches(4.45), Inches(4.1), Inches(2.4))
tf = box.text_frame
tf.word_wrap = True
tf.margin_left = tf.margin_right = Emu(0)
for i, t in enumerate(items_l):
    bullet(tf, t, size=12, color=TX, first=(i == 0))

# 오른쪽 패널 (terminal)
add_card(s2, Inches(5.2), Inches(3.5), Inches(7.7), Inches(3.4), accent=GD)
add_text(s2, Inches(5.45), Inches(3.65), Inches(6), Inches(0.4),
         "오른쪽 (65%)", size=11, bold=True, color=GD)
add_text(s2, Inches(5.45), Inches(3.95), Inches(6), Inches(0.4),
         "터미널 (claude 가 도는 곳)", size=15, bold=True, color=WHITE)
term_box = s2.shapes.add_textbox(Inches(5.45), Inches(4.45), Inches(7.2), Inches(1.8))
tf = term_box.text_frame
tf.word_wrap = True
tf.margin_left = tf.margin_right = Emu(0)
term_lines = [
    "user01@workshop:~$ claude",
    "> PLAN.md 만들어줘",
    "✓ Created PLAN.md",
]
for i, l in enumerate(term_lines):
    bullet(tf, l, size=12, color=GD, mono=True, first=(i == 0))

# AI 조수 칩 (오른쪽 하단)
add_rect(s2, Inches(11.0), Inches(6.05), Inches(1.85), Inches(0.7),
         fill=NV, line=BLM, line_w=Pt(1.2), radius=True)
add_text(s2, Inches(11.1), Inches(6.18), Inches(1.7), Inches(0.45),
         "🤖 AI 조수", size=12, bold=True, color=BLM, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

add_text(s2, Inches(0.6), Inches(7.05), Inches(12), Inches(0.4),
         "💡 진행 트래커는 하단에 — 1/8 ─ 2/8 ─ … ─ 8/8 로 미션 진척도 표시",
         size=11, color=TX3)

# ═══════════════════════════════════════════════════════
# Slide 3 — 왼쪽 안내 패널: GOAL/INPUT/OUTPUT
# ═══════════════════════════════════════════════════════
s3 = prs.slides.add_slide(BLANK)
add_bg(s3)
add_header(s3, "BRIEFING", "왼쪽 안내 — 무엇을 입력하고 무엇이 나와야 하나")

cards = [
    ("🎯 GOAL", GD, "이 미션의 목적",
     '예) "Plan 모드를 써서 보고서 워크플로우 4~5단계 설계 + PLAN.md 저장"'),
    ("📥 INPUT", OR, "내가 무엇을 입력하는가",
     '예) "claude 실행 후 /plan 슬래시 커맨드로 시작"'),
    ("📤 OUTPUT", BLM, "어떤 산출물이 나와야 하는가  +  📄 파일 칩",
     '파일 생성되면 회색→녹색 + 글로우 효과'),
]
y = Inches(2.3)
for i, (title, color, sub, ex) in enumerate(cards):
    cy = y + Inches(i * 1.55)
    add_card(s3, Inches(0.6), cy, Inches(12), Inches(1.35), accent=color)
    add_text(s3, Inches(0.95), cy + Inches(0.18), Inches(4), Inches(0.4),
             title, size=18, bold=True, color=color)
    add_text(s3, Inches(0.95), cy + Inches(0.6), Inches(11), Inches(0.4),
             sub, size=14, color=TX)
    add_text(s3, Inches(0.95), cy + Inches(0.9), Inches(11), Inches(0.4),
             ex, size=12, color=TX2, font=MONO)

# ═══════════════════════════════════════════════════════
# Slide 4 — 4 버튼
# ═══════════════════════════════════════════════════════
s4 = prs.slides.add_slide(BLANK)
add_bg(s4)
add_header(s4, "BUTTONS", "주요 버튼 4개")
add_text(s4, Inches(0.6), Inches(2.1), Inches(12), Inches(0.5),
         "프롬프트 → 터미널 → 파일 생성 → 채점 — 각 단계마다 한 버튼.",
         size=14, color=TX2)

btns = [
    ("📋", "복사", OR, "프롬프트를 클립보드로", "터미널에 직접 붙여넣고 싶을 때"),
    ("▶", "실행", GD_D, "프롬프트를 터미널에 자동 입력", "(Enter는 안 누름 — 검토 후 직접)"),
    ("🧹", "클리어", TX3, "터미널 화면 비우기", "화면이 지저분해질 때"),
    ("🤖", "AI 채점", BLM, "산출물을 AI 가 상세 평가", "결과 파일이 만들어진 후 (회색→파랑)"),
]
for i, (icon, name, color, desc, when) in enumerate(btns):
    cx = Inches(0.6 + (i * 3.15))
    cy = Inches(3.0)
    add_card(s4, cx, cy, Inches(2.95), Inches(3.7), accent=color)
    add_text(s4, cx + Inches(0.25), cy + Inches(0.3), Inches(2.7), Inches(0.7),
             icon, size=32, color=color, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_text(s4, cx + Inches(0.25), cy + Inches(1.3), Inches(2.7), Inches(0.5),
             name, size=18, bold=True, color=color, align=PP_ALIGN.CENTER)
    add_text(s4, cx + Inches(0.25), cy + Inches(1.95), Inches(2.7), Inches(0.7),
             desc, size=12, color=TX, align=PP_ALIGN.CENTER)
    add_text(s4, cx + Inches(0.25), cy + Inches(2.75), Inches(2.7), Inches(0.7),
             when, size=11, color=TX3, align=PP_ALIGN.CENTER)

# ═══════════════════════════════════════════════════════
# Slide 5 — 미션 진행 흐름 (6 step)
# ═══════════════════════════════════════════════════════
s5 = prs.slides.add_slide(BLANK)
add_bg(s5)
add_header(s5, "FLOW", "미션 진행 흐름 — 6 단계")

flow = [
    ("1", "안내 읽기", "왼쪽 GOAL / INPUT / OUTPUT", OR),
    ("2", "프롬프트 실행", "📋 복사 또는 ▶ 실행 버튼", GD_D),
    ("3", "AI 작업", "터미널에서 claude 가 파일 생성", BL),
    ("4", "자동 감지", "3초 폴링 → 📄 토스트 + 🎉 컨페티", BLM),
    ("5", "다음 슬라이드", "체크 통과 → 2초 후 자동 진행", GD),
    ("6", "(선택) AI 채점", "더 자세한 피드백 원할 때", OR_D),
]
y = Inches(2.45)
for i, (step, title, desc, color) in enumerate(flow):
    row, col = divmod(i, 3)
    cx = Inches(0.6 + col * 4.3)
    cy = y + Inches(row * 2.35)
    add_card(s5, cx, cy, Inches(4.1), Inches(2.05), accent=color)
    # step 동그라미
    circle = s5.shapes.add_shape(MSO_SHAPE.OVAL,
                                  cx + Inches(0.25), cy + Inches(0.3),
                                  Inches(0.65), Inches(0.65))
    circle.fill.solid()
    circle.fill.fore_color.rgb = color
    circle.line.fill.background()
    add_text(s5, cx + Inches(0.25), cy + Inches(0.3), Inches(0.65), Inches(0.65),
             step, size=20, bold=True, color=WHITE,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_text(s5, cx + Inches(1.05), cy + Inches(0.35), Inches(2.9), Inches(0.5),
             title, size=15, bold=True, color=WHITE)
    add_text(s5, cx + Inches(0.3), cy + Inches(1.1), Inches(3.7), Inches(0.85),
             desc, size=12, color=TX2)

add_text(s5, Inches(0.6), Inches(7.05), Inches(12), Inches(0.4),
         "💡 핵심: 자동 감지로 다음 슬라이드까지 자동 진행. AI 채점은 옵션.",
         size=12, color=GD)

# ═══════════════════════════════════════════════════════
# Slide 6 — 필수/도전 체크리스트
# ═══════════════════════════════════════════════════════
s6 = prs.slides.add_slide(BLANK)
add_bg(s6)
add_header(s6, "CHECKLIST", "필수 + 도전 — 두 단계 체크")

card_with_bullets(
    s6, Inches(0.6), Inches(2.4), Inches(6.1), Inches(4.6), GD,
    "✅ 필수 체크리스트 (Mandatory)",
    [
        "통과해야 미션 클리어",
        "예: PLAN.md 파일이 생성되었다",
        "예: 단계가 3개 이상 있다",
        "예: 한국어로 설명되어 있다",
        "",
        "▸ 파일 생성 시 자동 통과 표시",
        "▸ AI 채점 받으면 항목별 ✓/✗ + 코멘트",
    ],
    title_size=15,
)
card_with_bullets(
    s6, Inches(6.85), Inches(2.4), Inches(6.0), Inches(4.6), OR,
    "⭐ 도전 체크리스트 (Challenge)",
    [
        "선택 — 더 좋은 결과물을 위한 추가 기준",
        "예: 입력/출력/산출물이 명시됨",
        "예: 마크다운 헤딩으로 구조화됨",
        "예: 최종 산출물이 명시됨",
        "",
        "▸ AI 채점 시에만 평가됨",
        "▸ 통과 못해도 미션 클리어에는 무관",
    ],
    title_size=15,
)

add_text(s6, Inches(0.6), Inches(7.05), Inches(12), Inches(0.4),
         "💡 자동 통과 후에도 'AI 채점 받기' 로 자세한 피드백을 받을 수 있어요",
         size=12, color=TX3)

# ═══════════════════════════════════════════════════════
# Slide 7 — AI 채점 동작 방식
# ═══════════════════════════════════════════════════════
s7 = prs.slides.add_slide(BLANK)
add_bg(s7)
add_header(s7, "AI GRADING", "AI 채점 — 동작 방식")

steps = [
    ("1", "버튼 클릭", "🤖 AI 채점 받기 (파일 생성 후 활성화)"),
    ("2", "파일 읽기", "outputFiles 경로에서 산출물 수집"),
    ("3", "프롬프트 구성", "goal + inputDesc + outputDesc 가 루브릭"),
    ("4", "Claude 호출", "학습자 본인 크리덴셜로 claude -p --json-schema"),
    ("5", "결과 파싱", "score · passed · summary · items[ ]"),
    ("6", "UI 반영", "각 체크 항목 옆에 ✓/✗ + 코멘트 인라인 표시"),
]
for i, (n, t, d) in enumerate(steps):
    row, col = divmod(i, 2)
    cx = Inches(0.6 + col * 6.35)
    cy = Inches(2.4 + row * 1.55)
    add_card(s7, cx, cy, Inches(6.1), Inches(1.35), accent=BLM)
    circle = s7.shapes.add_shape(MSO_SHAPE.OVAL,
                                  cx + Inches(0.25), cy + Inches(0.35),
                                  Inches(0.6), Inches(0.6))
    circle.fill.solid()
    circle.fill.fore_color.rgb = BLM
    circle.line.fill.background()
    add_text(s7, cx + Inches(0.25), cy + Inches(0.35), Inches(0.6), Inches(0.6),
             n, size=18, bold=True, color=WHITE,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_text(s7, cx + Inches(1.0), cy + Inches(0.18), Inches(5), Inches(0.5),
             t, size=15, bold=True, color=BLM)
    add_text(s7, cx + Inches(1.0), cy + Inches(0.7), Inches(5), Inches(0.5),
             d, size=12, color=TX2)

# ═══════════════════════════════════════════════════════
# Slide 8 — AI 채점 응답 예시
# ═══════════════════════════════════════════════════════
s8 = prs.slides.add_slide(BLANK)
add_bg(s8)
add_header(s8, "RESPONSE", "AI 채점 응답 예시")

# 점수 박스
add_card(s8, Inches(0.6), Inches(2.4), Inches(12), Inches(1.1), accent=GD)
add_text(s8, Inches(0.95), Inches(2.55), Inches(2), Inches(0.7),
         "80점", size=32, bold=True, color=GD)
add_chip(s8, Inches(2.45), Inches(2.75), "✓ PASS", fill=GD_D, size=12)
add_text(s8, Inches(4.0), Inches(2.7), Inches(8.5), Inches(0.5),
         '"잘 작성된 PLAN.md 입니다. 단계가 명확합니다."',
         size=14, color=TX)

# 필수 체크리스트 박스
add_card(s8, Inches(0.6), Inches(3.7), Inches(6.0), Inches(3.3), accent=GD)
add_text(s8, Inches(0.85), Inches(3.85), Inches(5.5), Inches(0.4),
         "✅ 필수 체크리스트", size=12, bold=True, color=GD)
items_pass = [
    ("✓", "PLAN.md 파일이 생성됨", "확인됨", GD),
    ("✓", "단계가 3개 이상", "5단계 명시", GD),
    ("✓", "한국어로 설명됨", "한국어 OK", GD),
]
for i, (mark, name, comment, color) in enumerate(items_pass):
    iy = Inches(4.3 + i * 0.85)
    add_text(s8, Inches(0.95), iy, Inches(0.3), Inches(0.4),
             mark, size=14, bold=True, color=color)
    add_text(s8, Inches(1.3), iy, Inches(5.0), Inches(0.4),
             name, size=12, color=TX)
    add_text(s8, Inches(1.3), iy + Inches(0.32), Inches(5.0), Inches(0.4),
             "💭 " + comment, size=10, color=TX3)

# 도전 체크리스트 박스
add_card(s8, Inches(6.85), Inches(3.7), Inches(6.0), Inches(3.3), accent=OR)
add_text(s8, Inches(7.1), Inches(3.85), Inches(5.5), Inches(0.4),
         "⭐ 도전 체크리스트", size=12, bold=True, color=OR)
items_chal = [
    ("✓", "입력/출력/산출물 명시", "단계마다 표시됨", OR),
    ("⨯", "마크다운 헤딩 사용", "## 헤딩 추가하면 더 좋음", BAD),
]
for i, (mark, name, comment, color) in enumerate(items_chal):
    iy = Inches(4.3 + i * 1.0)
    add_text(s8, Inches(7.2), iy, Inches(0.3), Inches(0.4),
             mark, size=14, bold=True, color=color)
    add_text(s8, Inches(7.55), iy, Inches(5.0), Inches(0.4),
             name, size=12, color=TX)
    add_text(s8, Inches(7.55), iy + Inches(0.32), Inches(5.0), Inches(0.55),
             "💭 " + comment, size=10, color=TX3)

# ═══════════════════════════════════════════════════════
# Slide 9 — AI 조수 (어시스턴트 오버레이)
# ═══════════════════════════════════════════════════════
s9 = prs.slides.add_slide(BLANK)
add_bg(s9)
add_header(s9, "ASSISTANT", "🤖 AI 조수 — 옆에 앉은 코치")

add_text(s9, Inches(0.6), Inches(2.1), Inches(12), Inches(0.6),
         "터미널 우측 하단에 자동으로 떠 있는 작은 카드. 학습자 진행을 30초마다 살펴보고 짧게 코칭.",
         size=14, color=TX2)

# 왼쪽: 입력
card_with_bullets(
    s9, Inches(0.6), Inches(3.0), Inches(6.1), Inches(4.0), BLM,
    "📥 무엇을 보고 코칭하나",
    [
        "터미널 스크롤백 (최근 1000줄)",
        "산출 파일 내용 (있으면)",
        "미션 goal + 필수 체크리스트",
        "",
        "▸ 30초마다 자동 호출",
        "▸ 미션 진입 8초 후 첫 호출",
        "▸ '지금 확인' 버튼으로 즉시도 가능",
    ],
)
# 오른쪽: 응답
status_card_y = Inches(3.0)
add_card(s9, Inches(6.85), status_card_y, Inches(6.0), Inches(4.0), accent=GD)
add_text(s9, Inches(7.1), status_card_y + Inches(0.18), Inches(5.5), Inches(0.4),
         "📤 4가지 상태", size=15, bold=True, color=GD)

statuses = [
    ("👍", "good", "잘 진행 중", GD),
    ("✨", "almost", "거의 다 됨", WN),
    ("🤔", "stuck", "막힌 듯", OR),
    ("💡", "empty", "아직 시작 안 함", BLM),
]
for i, (icon, status, label, color) in enumerate(statuses):
    iy = status_card_y + Inches(0.65 + i * 0.78)
    add_text(s9, Inches(7.1), iy, Inches(0.5), Inches(0.5),
             icon, size=20, color=color)
    add_text(s9, Inches(7.7), iy + Inches(0.05), Inches(2), Inches(0.4),
             status, size=13, bold=True, color=color, font=MONO)
    add_text(s9, Inches(9.7), iy + Inches(0.05), Inches(3), Inches(0.4),
             label, size=13, color=TX)

# ═══════════════════════════════════════════════════════
# Slide 10 — 파일 생성 가시화
# ═══════════════════════════════════════════════════════
s10 = prs.slides.add_slide(BLANK)
add_bg(s10)
add_header(s10, "FEEDBACK", "파일 생성 가시화 — 결과를 즉시 인지")

add_text(s10, Inches(0.6), Inches(2.1), Inches(12), Inches(0.6),
         "claude 가 파일을 만든 그 순간, 학습자가 시각적으로 명확히 알 수 있게 — 3단 피드백.",
         size=14, color=TX2)

# 토스트 (큰 박스 위)
toast_y = Inches(3.0)
add_rect(s10, Inches(3.5), toast_y, Inches(6.3), Inches(0.95),
         fill=GD_D, line=GD, line_w=Pt(1.5), radius=True)
add_text(s10, Inches(3.85), toast_y + Inches(0.1), Inches(0.7), Inches(0.7),
         "📄", size=28, color=WHITE, anchor=MSO_ANCHOR.MIDDLE)
add_text(s10, Inches(4.6), toast_y + Inches(0.13), Inches(5), Inches(0.4),
         "파일 생성 완료", size=10, bold=True, color=WHITE)
add_text(s10, Inches(4.6), toast_y + Inches(0.42), Inches(5), Inches(0.45),
         "PLAN.md", size=15, bold=True, color=WHITE, font=MONO)

# 3 단 카드
items = [
    ("①", "📄 토스트 (4.5초)",
     "터미널 위 큰 토스트 — '파일 생성 완료 PLAN.md'\n팝업 애니메이션",
     GD_D),
    ("②", "파일 칩 색상 전환",
     "OUTPUT 섹션의 파일 칩이\n회색 → 녹색 + 글로우 효과",
     GD),
    ("③", "🎉 컨페티 + 자동 진행",
     "체크 통과 → 컨페티 → 2초 후\n다음 슬라이드로 자동 이동",
     OR),
]
for i, (n, title, desc, color) in enumerate(items):
    cx = Inches(0.6 + i * 4.3)
    cy = Inches(4.5)
    add_card(s10, cx, cy, Inches(4.1), Inches(2.4), accent=color)
    add_text(s10, cx + Inches(0.3), cy + Inches(0.2), Inches(0.5), Inches(0.5),
             n, size=22, bold=True, color=color)
    add_text(s10, cx + Inches(1.0), cy + Inches(0.25), Inches(3), Inches(0.5),
             title, size=14, bold=True, color=color)
    add_text(s10, cx + Inches(0.3), cy + Inches(1.05), Inches(3.7), Inches(1.3),
             desc, size=12, color=TX2)

# ═══════════════════════════════════════════════════════
# Slide 11 — 시연 슬라이드 Haiku 자동 전환
# ═══════════════════════════════════════════════════════
s11 = prs.slides.add_slide(BLANK)
add_bg(s11)
add_header(s11, "DEMO MODE", "시연 슬라이드 Haiku 자동 전환 — 강사용")

add_text(s11, Inches(0.6), Inches(2.1), Inches(12), Inches(0.7),
         "시연 중 Claude 가 너무 길게 답해서 지루해지는 문제 해결.\n'시연:' 으로 시작하는 슬라이드 진입 시 모든 사용자가 자동으로 Haiku 4.5 로 전환.",
         size=14, color=TX2)

# 두 모드 비교
card_with_bullets(
    s11, Inches(0.6), Inches(3.5), Inches(6.1), Inches(3.5), OR,
    "⚡ 시연 모드 (Haiku 4.5)",
    [
        "슬라이드 제목이 '시연:' 으로 시작할 때",
        "응답 토큰 1024 (짧고 빠름)",
        "MAX_THINKING_TOKENS=0",
        "",
        "▸ 데모 시간 단축",
        "▸ 강사는 추가로 터미널 자동 클리어",
        "▸ 학습자도 swap 됨 (다음 claude 호출부터)",
    ],
)
card_with_bullets(
    s11, Inches(6.85), Inches(3.5), Inches(6.0), Inches(3.5), GD,
    "🎯 일반 / 미션 모드 (Sonnet 4.6)",
    [
        "그 외 모든 슬라이드에서",
        "응답 토큰 8192 (풀 품질)",
        "추론 토큰 충분히 사용",
        "",
        "▸ 미션 진입 시 자동으로 normal 강제",
        "▸ MissionSlide 가 setDemoMode('normal') 호출",
        "▸ 학습자가 본인 작업할 땐 항상 풀 품질",
    ],
)

# ═══════════════════════════════════════════════════════
# Slide 12 — 트러블슈팅
# ═══════════════════════════════════════════════════════
s12 = prs.slides.add_slide(BLANK)
add_bg(s12)
add_header(s12, "TROUBLESHOOTING", "자주 겪는 문제 + 해결")

issues = [
    ("AI 채점 버튼이 회색이에요",
     "결과 파일이 아직 안 만들어졌기 때문. 먼저 프롬프트를 실행해서 파일을 생성하세요.",
     OR),
    ("'채점 실패' 로 나와요",
     "1) 학습자 크리덴셜 만료 → claude /login 다시.  2) 파일이 너무 큼 (64KB 초과)",
     BAD),
    ("AI 조수가 계속 stuck 으로만 떠요",
     "터미널에 거의 입력이 없을 때 정상 동작. 닫기/최소화 가능.",
     WN),
    ("파일 생성 토스트가 안 떠요",
     "false→true 전환 순간만 표시. 이미 통과한 미션을 다시 보면 안 뜸 (의도된 동작).",
     BLM),
    ("시연 Haiku 가 적용 안 돼요",
     "슬라이드 제목이 '시연:' 로 시작해야 함. 다른 접두어는 무시.",
     OR_D),
]
for i, (q, a, color) in enumerate(issues):
    cy = Inches(2.3 + i * 0.95)
    add_card(s12, Inches(0.6), cy, Inches(12), Inches(0.85), accent=color)
    add_text(s12, Inches(0.95), cy + Inches(0.13), Inches(11), Inches(0.4),
             "Q. " + q, size=13, bold=True, color=color)
    add_text(s12, Inches(0.95), cy + Inches(0.45), Inches(11), Inches(0.4),
             "A. " + a, size=12, color=TX2)

# ═══════════════════════════════════════════════════════
# Slide 13 — 자동 테스트 5 스위트
# ═══════════════════════════════════════════════════════
s13 = prs.slides.add_slide(BLANK)
add_bg(s13)
add_header(s13, "TESTS", "자동 테스트 — 5 스위트, 124+ 단언")

add_text(s13, Inches(0.6), Inches(2.1), Inches(12), Inches(0.5),
         "본 워크북의 새 기능은 자동 테스트로 검증됩니다. bash tests/run-all.sh 로 한 번에 실행.",
         size=14, color=TX2)

suites = [
    ("[1/5]", "셸 스크립트 단위", "admin-action / coach / grade-mission 인자 검증", "17", GD),
    ("[2/5]", "server.js 라우트 통합", "/api/health, /me, /my-scrollback, /coach, /grade-mission", "23", BLM),
    ("[3/5]", "프론트엔드 정적 + 빌드", "신규 컴포넌트 + runtime helpers + vite build", "39", OR),
    ("[4/5]", "UI 렌더 + preview", "레이아웃 토큰, 안티패턴, vite preview 스모크", "20", BL),
    ("[5/5]", "학습자 여정 (user01)", "인증→체크→파일→채점→코칭→스크롤백→세션초기화", "25", GD_D),
]
for i, (idx, title, desc, count, color) in enumerate(suites):
    cy = Inches(2.95 + i * 0.85)
    add_card(s13, Inches(0.6), cy, Inches(12), Inches(0.75), accent=color)
    add_text(s13, Inches(0.9), cy + Inches(0.18), Inches(1), Inches(0.4),
             idx, size=11, bold=True, color=color, font=MONO)
    add_text(s13, Inches(1.85), cy + Inches(0.13), Inches(4), Inches(0.4),
             title, size=14, bold=True, color=WHITE)
    add_text(s13, Inches(1.85), cy + Inches(0.42), Inches(8), Inches(0.4),
             desc, size=11, color=TX2)
    add_chip(s13, Inches(11.2), cy + Inches(0.2), count + " 단언", fill=color, fg=WHITE, size=11)

add_text(s13, Inches(0.6), Inches(7.2), Inches(12), Inches(0.3),
         "✓ 124 / 124 통과 — 5 스위트 모두 grün",
         size=12, color=GD, font=MONO, align=PP_ALIGN.CENTER)

# ═══════════════════════════════════════════════════════
# Slide 14 — 종료 / 변경 파일 요약
# ═══════════════════════════════════════════════════════
s14 = prs.slides.add_slide(BLANK)
add_bg(s14)
add_rect(s14, 0, 0, SLIDE_W, SLIDE_H, fill=BG)
add_rect(s14, 0, 0, Inches(5), SLIDE_H, fill=NV)
add_text(s14, Inches(0.6), Inches(0.6), Inches(4), Inches(0.5),
         "MIRAE ASSET LIFE", size=14, bold=True, color=WHITE)
add_text(s14, Inches(0.6), Inches(1.0), Inches(4), Inches(0.4),
         "Workbook UI Guide · End", size=12, color=TX2)
add_text(s14, Inches(0.6), Inches(2.5), Inches(4), Inches(2),
         "끝!", size=72, bold=True, color=OR)
add_text(s14, Inches(0.6), Inches(4.8), Inches(4), Inches(0.5),
         "더 알고 싶다면", size=12, color=TX3)
add_text(s14, Inches(0.6), Inches(5.2), Inches(4), Inches(0.4),
         "docs/handout-workbook-ui.md", size=12, color=OR, font=MONO)
add_text(s14, Inches(0.6), Inches(5.55), Inches(4), Inches(0.4),
         "tests/run-all.sh", size=12, color=OR, font=MONO)

# 우측 본문
add_text(s14, Inches(5.6), Inches(0.8), Inches(7.4), Inches(0.5),
         "CHANGED FILES", size=11, bold=True, color=OR)
add_text(s14, Inches(5.6), Inches(1.2), Inches(7.4), Inches(0.6),
         "변경된 파일 (개발자용)", size=22, bold=True, color=WHITE)

groups = [
    ("프론트엔드", OR, [
        "src/components/MissionSlide.jsx — 안내/채점/토스트",
        "src/components/AssistantOverlay.jsx (신규)",
        "src/lib/runtime.js — coachMission · fetchMyScrollback",
        "src/workbook.jsx — 워크북 사용법 슬라이드 + Haiku",
    ]),
    ("백엔드", BLM, [
        "deploy/multi-user/api-server/server.js — /coach · /scrollback",
        "deploy/multi-user/api-server/admin-action.sh — capture-scrollback",
        "deploy/multi-user/api-server/coach.sh (신규)",
        "deploy/multi-user/setup-server.sh — sudoers",
    ]),
    ("테스트", GD, [
        "tests/run-all.sh + 5개 스위트",
        "tests/fixtures/mock-*.sh (4개 모킹)",
        "총 124+ 단언, 모두 통과",
    ]),
]
y = Inches(2.1)
for title, color, items in groups:
    add_text(s14, Inches(5.6), y, Inches(7.4), Inches(0.4),
             title, size=13, bold=True, color=color)
    for i, item in enumerate(items):
        add_text(s14, Inches(5.7), y + Inches(0.4 + i * 0.32),
                 Inches(7.4), Inches(0.32),
                 "  ▸ " + item, size=11, color=TX2, font=MONO)
    y += Inches(0.45 + len(items) * 0.32 + 0.1)

# ═══════════════════════════════════════════════════════
# 저장
# ═══════════════════════════════════════════════════════
prs.save(OUT)
print(f"✓ {OUT} ({len(prs.slides)} slides)")
