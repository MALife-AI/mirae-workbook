#!/usr/bin/env python3
"""집에서 Claude Code 실습하는 방법 — 워크숍 종료 후 배포용 별도 PPTX.

handout.md Appendix A 내용을 압축해서 14장 슬라이드로.
브랜드: 미래에셋생명 · 오렌지 #F58220 + 짙은 네이비.
"""
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

DOCS = Path(__file__).parent
OUT = DOCS / "home-practice.pptx"

# ── 브랜드 색상 ────────────────────────────────────────
OR = RGBColor(0xF5, 0x82, 0x20)
OR_D = RGBColor(0xCB, 0x60, 0x15)
NV = RGBColor(0x04, 0x3B, 0x72)
NV_D = RGBColor(0x06, 0x1E, 0x30)
BG = RGBColor(0x04, 0x18, 0x28)
TX = RGBColor(0xE5, 0xE8, 0xEC)
TX2 = RGBColor(0x8D, 0xA0, 0xB8)
TX3 = RGBColor(0x5A, 0x7A, 0x98)
GD = RGBColor(0x86, 0xEF, 0xAC)
BAD = RGBColor(0xFC, 0xA5, 0xA5)
WN = RGBColor(0xFB, 0xBF, 0x24)
BD = RGBColor(0x0A, 0x30, 0x50)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BLACK = RGBColor(0x00, 0x00, 0x00)

KO_FONT = "Pretendard"
MONO = "JetBrains Mono"

# ── 슬라이드 사이즈: 16:9 와이드 ───────────────────────
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

SLIDE_W = prs.slide_width
SLIDE_H = prs.slide_height
BLANK = prs.slide_layouts[6]  # 빈 레이아웃


# ── 헬퍼 ──────────────────────────────────────────────
def add_bg(slide, color=BG):
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H)
    bg.fill.solid()
    bg.fill.fore_color.rgb = color
    bg.line.fill.background()
    return bg


def add_text(slide, x, y, w, h, text, *, size=18, color=TX, bold=False,
             font=KO_FONT, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP):
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Emu(0)
    tf.margin_top = tf.margin_bottom = Emu(0)
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return box


def add_rect(slide, x, y, w, h, *, fill=None, line=None, line_w=None, radius=False):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE,
        x, y, w, h,
    )
    if fill is not None:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    else:
        shape.fill.background()
    if line is not None:
        shape.line.color.rgb = line
        if line_w:
            shape.line.width = line_w
    else:
        shape.line.fill.background()
    return shape


def add_chip(slide, x, y, text, *, fill=OR, fg=WHITE, size=11, padx=8, pady=4):
    """작은 chip — 텍스트 길이에 맞춰서 자동 너비."""
    w = Emu(int(len(text) * 110000 + padx * 18000))
    h = Inches(0.32)
    chip = add_rect(slide, x, y, w, h, fill=fill, radius=True)
    chip.line.fill.background()
    tf = chip.text_frame
    tf.margin_left = tf.margin_right = Emu(60000)
    tf.margin_top = tf.margin_bottom = Emu(0)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = text
    r.font.name = KO_FONT
    r.font.size = Pt(size)
    r.font.bold = True
    r.font.color.rgb = fg
    return chip


def add_header(slide, kicker, title):
    """모든 본문 슬라이드 공통 헤더 — 좌상단 kicker chip + 큰 제목."""
    add_chip(slide, Inches(0.6), Inches(0.5), kicker)
    add_text(
        slide, Inches(0.6), Inches(0.95), Inches(12), Inches(0.9),
        title, size=32, bold=True, color=WHITE,
    )
    # 오렌지 언더라인
    line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0.6), Inches(1.85), Inches(0.7), Emu(40000),
    )
    line.fill.solid()
    line.fill.fore_color.rgb = OR
    line.line.fill.background()
    # footer 페이지 번호
    add_text(
        slide, Inches(12.5), Inches(7.05), Inches(0.7), Inches(0.3),
        f"{len(prs.slides) + 1:02d}",
        size=10, color=TX3, font=MONO, align=PP_ALIGN.RIGHT,
    )


def add_card(slide, x, y, w, h, *, accent=OR, fill_alpha=True):
    """오렌지 좌측 액센트 + 카드 박스."""
    # 본체
    card = add_rect(slide, x, y, w, h, fill=NV_D, line=BD, line_w=Pt(0.75), radius=True)
    # 좌측 액센트 바
    bar = add_rect(slide, x, y, Inches(0.06), h, fill=accent)
    return card


def bullet(tf, text, *, size=13, color=TX, bold=False, indent=0, mono=False, first=False):
    p = tf.paragraphs[0] if first else tf.add_paragraph()
    p.alignment = PP_ALIGN.LEFT
    p.level = indent
    p.space_after = Pt(3)
    r = p.add_run()
    r.text = text
    r.font.name = MONO if mono else KO_FONT
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.color.rgb = color
    return p


# ═══════════════════════════════════════════════════════
# Slide 1 — 표지
# ═══════════════════════════════════════════════════════
s1 = prs.slides.add_slide(BLANK)
add_bg(s1)
# 큰 오렌지 블록 (왼쪽)
add_rect(s1, 0, 0, Inches(5), SLIDE_H, fill=OR_D)
add_text(
    s1, Inches(0.6), Inches(0.6), Inches(4), Inches(0.5),
    "MIRAE ASSET LIFE", size=14, bold=True, color=WHITE,
)
add_text(
    s1, Inches(0.6), Inches(1.0), Inches(4), Inches(0.4),
    "Vibe Coding Workshop", size=12, color=WHITE,
)

# 우측 본문
add_text(
    s1, Inches(5.6), Inches(1.7), Inches(7.4), Inches(0.5),
    "AFTER THE WORKSHOP", size=12, bold=True, color=OR,
)
add_text(
    s1, Inches(5.6), Inches(2.2), Inches(7.4), Inches(2.0),
    "집에서\nClaude Code\n실습하기",
    size=44, bold=True, color=WHITE,
)
# 부제
add_text(
    s1, Inches(5.6), Inches(5.0), Inches(7.4), Inches(0.6),
    "워크숍이 끝난 뒤,\n본인 PC 에서 직접 사용하는 방법",
    size=18, color=TX2,
)
# 하단 메타
add_text(
    s1, Inches(5.6), Inches(6.6), Inches(7.4), Inches(0.4),
    "미래에셋생명  ·  2026",
    size=11, color=TX3, font=MONO,
)

# ═══════════════════════════════════════════════════════
# Slide 2 — 개요: 두 갈래
# ═══════════════════════════════════════════════════════
s2 = prs.slides.add_slide(BLANK)
add_bg(s2)
add_header(s2, "OVERVIEW", "Claude 를 집에서 쓰는 두 가지 방법")

# 본문 인트로
add_text(
    s2, Inches(0.6), Inches(2.1), Inches(12), Inches(0.6),
    "워크숍에서 사용한 환경은 5시간 동안만 유지됩니다.\n이후엔 본인 PC 에 직접 설치해서 계속 사용하세요.",
    size=15, color=TX2,
)

# 두 카드
def two_card(x, accent, title, lines):
    add_card(s2, x, Inches(3.6), Inches(6.1), Inches(3.4), accent=accent)
    add_text(s2, x + Inches(0.3), Inches(3.7), Inches(5.8), Inches(0.4),
             title, size=18, bold=True, color=accent)
    box = s2.shapes.add_textbox(x + Inches(0.3), Inches(4.2), Inches(5.6), Inches(2.7))
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Emu(0)
    for i, ln in enumerate(lines):
        bullet(tf, ln, size=13, color=TX, first=(i == 0))


two_card(
    Inches(0.6), OR,
    "⚡ Claude Code (CLI)",
    [
        "터미널 기반 — 워크숍에서 쓴 그것",
        "파일 자동 생성·수정·실행 가능",
        "Skill / Command / Hook 으로 자동화",
        "보고서 · 분석 · 코드 작업에 강력",
        "약간의 터미널 익숙함 필요",
    ],
)
two_card(
    Inches(6.7), RGBColor(0x38, 0xBD, 0xF8),
    "🖥️ Claude Desktop (앱)",
    [
        "GUI 채팅 앱 — 카톡처럼 친숙",
        "drag & drop 으로 파일 첨부",
        "비개발자도 즉시 사용 가능",
        "글쓰기 · 일상 질문 · 아이디어",
        "파일 생성 · 자동화는 제한적",
    ],
)

add_text(
    s2, Inches(0.6), Inches(7.05), Inches(12), Inches(0.4),
    "💡 처음엔 둘 다 설치해 두고 상황에 맞게 골라 쓰는 것을 권장합니다.",
    size=12, color=GD,
)

# ═══════════════════════════════════════════════════════
# Slide 3 — Code 사전 준비
# ═══════════════════════════════════════════════════════
s3 = prs.slides.add_slide(BLANK)
add_bg(s3)
add_header(s3, "STEP 1 / CLI", "사전 준비 — 필요한 것")

cards = [
    ("Node.js 18 이상", "node --version", "nodejs.org 에서 LTS 다운로드"),
    ("Git", "git --version", "Windows: Git for Windows (Git Bash 포함)"),
    ("Anthropic 계정", "claude.ai 가입", "Pro / Max 구독 권장 — 정액제"),
]
y = Inches(2.3)
for i, (title, cmd, desc) in enumerate(cards):
    cy = y + Inches(i * 1.55)
    add_card(s3, Inches(0.6), cy, Inches(12), Inches(1.35))
    add_text(s3, Inches(0.9), cy + Inches(0.15), Inches(4), Inches(0.4),
             title, size=16, bold=True, color=OR)
    add_text(s3, Inches(0.9), cy + Inches(0.6), Inches(5), Inches(0.4),
             cmd, size=13, color=TX, font=MONO)
    add_text(s3, Inches(6.5), cy + Inches(0.4), Inches(6), Inches(0.5),
             desc, size=13, color=TX2)

# ═══════════════════════════════════════════════════════
# Slide 4 — Code 설치
# ═══════════════════════════════════════════════════════
s4 = prs.slides.add_slide(BLANK)
add_bg(s4)
add_header(s4, "STEP 2 / CLI", "Claude Code 설치 — 한 줄")

add_text(
    s4, Inches(0.6), Inches(2.2), Inches(12), Inches(0.5),
    "OS 와 무관 — Mac · Windows · Linux 모두 동일:",
    size=15, color=TX2,
)

# 큰 코드 블록
add_card(s4, Inches(0.6), Inches(3.0), Inches(12), Inches(1.4), accent=OR)
add_text(
    s4, Inches(1.0), Inches(3.45), Inches(11), Inches(0.6),
    "npm install -g @anthropic-ai/claude-code",
    size=24, bold=True, color=OR, font=MONO,
)

add_text(
    s4, Inches(0.6), Inches(4.7), Inches(12), Inches(0.5),
    "✅ 설치가 끝나면 어디서든 claude 한 줄로 시작:",
    size=14, color=GD,
)

add_card(s4, Inches(0.6), Inches(5.3), Inches(12), Inches(1.6), accent=GD)
box = s4.shapes.add_textbox(Inches(1.0), Inches(5.45), Inches(11), Inches(1.4))
tf = box.text_frame
tf.word_wrap = True
bullet(tf, "$ cd ~/내-프로젝트-폴더", size=14, color=TX, mono=True, first=True)
bullet(tf, "$ claude", size=14, color=OR, mono=True, bold=True)
bullet(tf, "  → Claude Code 가 시작되고 인증 안내가 나타납니다", size=12, color=TX3)

# ═══════════════════════════════════════════════════════
# Slide 5 — 첫 로그인
# ═══════════════════════════════════════════════════════
s5 = prs.slides.add_slide(BLANK)
add_bg(s5)
add_header(s5, "STEP 3 / CLI", "첫 로그인 — 한 번만 하면 영구")

steps = [
    ("1", "claude 입력", "터미널에서 그대로", "claude"),
    ("2", "/login 입력", "Claude 안에서", "/login"),
    ("3", "URL 열기", "표시된 URL 을 본인 PC 브라우저에", "https://claude.ai/oauth/..."),
    ("4", "코드 복사·붙여넣기", "받은 인증 코드를 터미널에", "(자동 인증 완료)"),
]
y = Inches(2.3)
for i, (n, title, desc, cmd) in enumerate(steps):
    cy = y + Inches(i * 1.1)
    # 번호 원
    circle = s5.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.6), cy + Inches(0.18), Inches(0.7), Inches(0.7))
    circle.fill.solid(); circle.fill.fore_color.rgb = OR
    circle.line.fill.background()
    tf = circle.text_frame
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = n; r.font.size = Pt(20); r.font.bold = True; r.font.color.rgb = WHITE; r.font.name = KO_FONT
    # 본문
    add_text(s5, Inches(1.5), cy + Inches(0.15), Inches(4), Inches(0.5),
             title, size=16, bold=True, color=WHITE)
    add_text(s5, Inches(1.5), cy + Inches(0.6), Inches(4), Inches(0.4),
             desc, size=12, color=TX2)
    add_text(s5, Inches(6.5), cy + Inches(0.35), Inches(6.5), Inches(0.5),
             cmd, size=14, color=OR, font=MONO)

add_text(
    s5, Inches(0.6), Inches(6.95), Inches(12), Inches(0.4),
    "💾 한 번 인증하면 ~/.claude/.credentials.json 에 저장 — 다음부터는 자동",
    size=12, color=GD,
)

# ═══════════════════════════════════════════════════════
# Slide 6 — 첫 명령
# ═══════════════════════════════════════════════════════
s6 = prs.slides.add_slide(BLANK)
add_bg(s6)
add_header(s6, "STEP 4 / CLI", "첫 명령 — 워크숍에서 한 그것 그대로")

add_text(
    s6, Inches(0.6), Inches(2.2), Inches(12), Inches(0.5),
    "한국어로 자연스럽게 지시하면 됩니다:",
    size=15, color=TX2,
)

prompts = [
    ("📊", "이 폴더의 *.csv 를 합쳐서 summary.xlsx 만들어줘"),
    ("📝", "퇴직연금 시장 동향 보고서 만들어서 outputs/ 에 저장해줘"),
    ("🔍", "PLAN.md 를 읽고 1단계 작업을 시작해줘"),
    ("🐛", "이 .py 파일에서 버그 찾아서 고쳐줘"),
]
y = Inches(2.95)
for i, (ic, p) in enumerate(prompts):
    cy = y + Inches(i * 0.85)
    add_card(s6, Inches(0.6), cy, Inches(12), Inches(0.7))
    add_text(s6, Inches(0.85), cy + Inches(0.18), Inches(0.5), Inches(0.4),
             ic, size=18, color=WHITE)
    add_text(s6, Inches(1.4), cy + Inches(0.18), Inches(11), Inches(0.4),
             p, size=14, color=TX, font=MONO)

add_text(
    s6, Inches(0.6), Inches(6.8), Inches(12), Inches(0.5),
    "✅ Claude 가 코드 작성 → 실행 → 결과 파일 생성. 사람은 결과만 검토.",
    size=12, color=GD,
)

# ═══════════════════════════════════════════════════════
# Slide 7 — Windows 사용자 주의
# ═══════════════════════════════════════════════════════
s7 = prs.slides.add_slide(BLANK)
add_bg(s7)
add_header(s7, "WINDOWS", "Windows 사용자 주의사항")

# 경고 카드
add_card(s7, Inches(0.6), Inches(2.3), Inches(12), Inches(2.2), accent=WN)
add_text(s7, Inches(0.9), Inches(2.45), Inches(11.5), Inches(0.5),
         "⚠️ Claude Code 는 내부적으로 bash 셸을 사용합니다",
         size=18, bold=True, color=WN)
box = s7.shapes.add_textbox(Inches(0.9), Inches(3.0), Inches(11.5), Inches(1.5))
tf = box.text_frame; tf.word_wrap = True
bullet(tf, "→ Git Bash 또는 WSL 환경이 필요합니다", size=14, color=TX, first=True)
bullet(tf, "→ Git for Windows 를 설치하면 Git Bash 가 함께 설치됩니다", size=14, color=TX2)
bullet(tf, "→ PowerShell / cmd 에서도 동작은 하지만 Git Bash 가 가장 안정적", size=14, color=TX2)

# 권장 단계
add_card(s7, Inches(0.6), Inches(4.7), Inches(12), Inches(2.4), accent=GD)
add_text(s7, Inches(0.9), Inches(4.85), Inches(11.5), Inches(0.5),
         "✅ 권장 설치 순서",
         size=16, bold=True, color=GD)
box = s7.shapes.add_textbox(Inches(0.9), Inches(5.4), Inches(11.5), Inches(1.7))
tf = box.text_frame; tf.word_wrap = True
bullet(tf, "1. gitforwindows.org → Git for Windows 설치 (Git Bash 포함)",
       size=13, color=TX, first=True)
bullet(tf, "2. nodejs.org → Node.js LTS 설치", size=13, color=TX)
bullet(tf, "3. Git Bash 열고: npm install -g @anthropic-ai/claude-code", size=13, color=TX, mono=True)
bullet(tf, "4. claude → /login → 끝", size=13, color=OR, mono=True, bold=True)

# ═══════════════════════════════════════════════════════
# Slide 8 — Desktop 설치
# ═══════════════════════════════════════════════════════
s8 = prs.slides.add_slide(BLANK)
add_bg(s8)
add_header(s8, "ALTERNATIVE", "Claude Desktop (GUI 앱) 설치")

add_text(s8, Inches(0.6), Inches(2.2), Inches(12), Inches(0.5),
         "claude.ai 의 다운로드 페이지에서 본인 OS 에 맞는 설치 파일을 받으세요:",
         size=15, color=TX2)

# OS 카드 3개
for i, (icon, os_name, fname, accent) in enumerate([
    ("🪟", "Windows", "Claude-Setup.exe", RGBColor(0x38, 0xBD, 0xF8)),
    ("🍎", "macOS", "Claude.dmg", RGBColor(0xA8, 0xC0, 0xDA)),
    ("🐧", "Linux", ".deb / AppImage", RGBColor(0x86, 0xEF, 0xAC)),
]):
    cx = Inches(0.6 + i * 4.13)
    add_card(s8, cx, Inches(3.0), Inches(4), Inches(2.0), accent=accent)
    add_text(s8, cx + Inches(0.3), Inches(3.15), Inches(3.5), Inches(0.6),
             icon + " " + os_name, size=20, bold=True, color=accent)
    add_text(s8, cx + Inches(0.3), Inches(3.85), Inches(3.5), Inches(0.5),
             fname, size=13, color=TX, font=MONO)
    add_text(s8, cx + Inches(0.3), Inches(4.35), Inches(3.5), Inches(0.5),
             "다운로드 → 다음·다음·완료", size=12, color=TX2)

# 첫 실행 안내
add_card(s8, Inches(0.6), Inches(5.4), Inches(12), Inches(1.6), accent=OR)
add_text(s8, Inches(0.9), Inches(5.55), Inches(11.5), Inches(0.5),
         "첫 실행", size=16, bold=True, color=OR)
box = s8.shapes.add_textbox(Inches(0.9), Inches(6.0), Inches(11.5), Inches(1.0))
tf = box.text_frame; tf.word_wrap = True
bullet(tf, "1. 앱 실행 → Claude 계정으로 로그인 → 채팅 화면", size=13, color=TX, first=True)
bullet(tf, '2. 바로 입력 가능: "오늘 해야 할 일을 정리해 줘"', size=13, color=TX)

# ═══════════════════════════════════════════════════════
# Slide 9 — CLI vs Desktop 비교 표
# ═══════════════════════════════════════════════════════
s9 = prs.slides.add_slide(BLANK)
add_bg(s9)
add_header(s9, "COMPARE", "CLI vs Desktop 한눈에")

rows = [
    ("구분", "Claude Code (CLI)", "Claude Desktop"),
    ("인터페이스", "터미널 (검은 창)", "채팅 앱 (윈도우/맥)"),
    ("학습 곡선", "약간 있음", "매우 쉬움 (카톡 수준)"),
    ("파일 생성·수정", "✓ 자유롭게", "△ 첨부만"),
    ("코드 직접 실행", "✓ Bash, Python, Node …", "✗ 불가"),
    ("폴더 단위 작업", "✓ 전체 프로젝트 분석", "△ 단일 첨부"),
    ("Skill / Hook / Command", "✓ 모두 사용", "✗ 사용 불가"),
    ("자동화 워크플로우", "✓ 핵심 기능", "△ 거의 불가"),
    ("적합한 용도", "보고서 자동화·코드", "일상 질문·글쓰기"),
]

# 표 그리기
table_x = Inches(0.6); table_y = Inches(2.2)
col_w = [Inches(2.6), Inches(5.0), Inches(4.4)]
row_h = Inches(0.5)
total_w = sum(col_w, Emu(0))
table_shape = s9.shapes.add_table(len(rows), 3, table_x, table_y, total_w, row_h * len(rows))
table = table_shape.table
for i, w in enumerate(col_w):
    table.columns[i].width = w
for ri in range(len(rows)):
    table.rows[ri].height = row_h
for ri, row in enumerate(rows):
    for ci, cell_text in enumerate(row):
        cell = table.cell(ri, ci)
        cell.fill.solid()
        if ri == 0:
            cell.fill.fore_color.rgb = OR
            color = WHITE; bold = True; size = 14
        else:
            cell.fill.fore_color.rgb = NV_D if ri % 2 == 1 else BG
            if ci == 1:
                color = OR; bold = True; size = 13
            elif ci == 2:
                color = TX2; bold = False; size = 13
            else:
                color = TX; bold = True; size = 13
        cell.margin_left = cell.margin_right = Emu(80000)
        cell.margin_top = cell.margin_bottom = Emu(40000)
        tf = cell.text_frame
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT if ci == 0 else PP_ALIGN.LEFT
        # 기본 run 클리어
        p.text = ""
        r = p.add_run()
        r.text = cell_text
        r.font.size = Pt(size); r.font.bold = bold; r.font.color.rgb = color; r.font.name = KO_FONT

add_text(s9, Inches(0.6), Inches(7.05), Inches(12), Inches(0.4),
         "→ 가벼운 질문 · 글쓰기 = Desktop  |  보고서 자동화 · 파일 작업 = Code",
         size=12, color=GD)

# ═══════════════════════════════════════════════════════
# Slide 10 — 자주 쓰는 슬래시 커맨드
# ═══════════════════════════════════════════════════════
s10 = prs.slides.add_slide(BLANK)
add_bg(s10)
add_header(s10, "REFERENCE", "자주 쓰는 슬래시 커맨드")

cmds = [
    ("/login", "Claude 계정 로그인"),
    ("/logout", "로그아웃"),
    ("/clear", "현재 대화 초기화"),
    ("/compact", "대화 요약 (긴 대화 시)"),
    ("/memory", "영구 메모리 관리"),
    ("/context", "컨텍스트 사용량 확인"),
    ("/usage", "Claude Max 사용량 확인"),
    ("/plan", "Plan 모드 — 작업 계획부터"),
    ("/help", "명령어 목록"),
    ("/exit", "Claude Code 종료"),
]
# 2 컬럼 5 줄
for i, (cmd, desc) in enumerate(cmds):
    col = i // 5; row = i % 5
    cx = Inches(0.6 + col * 6.2); cy = Inches(2.3 + row * 0.95)
    add_card(s10, cx, cy, Inches(6), Inches(0.8))
    add_text(s10, cx + Inches(0.2), cy + Inches(0.18), Inches(2.5), Inches(0.5),
             cmd, size=15, bold=True, color=OR, font=MONO)
    add_text(s10, cx + Inches(2.6), cy + Inches(0.22), Inches(3.4), Inches(0.5),
             desc, size=12, color=TX2)

# ═══════════════════════════════════════════════════════
# Slide 11 — 작업 흐름
# ═══════════════════════════════════════════════════════
s11 = prs.slides.add_slide(BLANK)
add_bg(s11)
add_header(s11, "WORKFLOW", "일반적인 작업 흐름")

steps = [
    ("1", "작업 폴더로 이동", "cd ~/projects/my-task"),
    ("2", "Claude Code 실행", "claude"),
    ("3", "한국어로 지시", '> 이 폴더의 *.csv 를 합쳐서 summary.xlsx 만들어줘'),
    ("4", "Claude 가 작업", "→ 코드 작성 · 실행 · 결과 파일 생성"),
    ("5", "필요시 수정 요청", '> 매출 컬럼은 천 단위 콤마 추가해줘'),
    ("6", "종료", "/exit  또는  Ctrl+C 두 번"),
]
for i, (n, title, cmd) in enumerate(steps):
    cy = Inches(2.3 + i * 0.78)
    # 번호 박스
    box = add_rect(s11, Inches(0.6), cy + Inches(0.05), Inches(0.55), Inches(0.55), fill=OR, radius=True)
    box.line.fill.background()
    tf = box.text_frame
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_top = tf.margin_bottom = Emu(0)
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = n; r.font.size = Pt(18); r.font.bold = True; r.font.color.rgb = WHITE; r.font.name = KO_FONT

    add_text(s11, Inches(1.4), cy + Inches(0.0), Inches(4), Inches(0.4),
             title, size=15, bold=True, color=WHITE)
    add_text(s11, Inches(1.4), cy + Inches(0.35), Inches(11.5), Inches(0.4),
             cmd, size=12, color=OR, font=MONO)

# ═══════════════════════════════════════════════════════
# Slide 12 — 안전하게 쓰기 (권한, 위험 명령)
# ═══════════════════════════════════════════════════════
s12 = prs.slides.add_slide(BLANK)
add_bg(s12)
add_header(s12, "SAFETY", "안전하게 쓰기 — 꼭 지킬 것")

# 좌측 권한 카드
add_card(s12, Inches(0.6), Inches(2.3), Inches(6.1), Inches(4.6), accent=OR)
add_text(s12, Inches(0.9), Inches(2.45), Inches(5.5), Inches(0.5),
         "권한과 폴더", size=18, bold=True, color=OR)
box = s12.shapes.add_textbox(Inches(0.9), Inches(2.95), Inches(5.7), Inches(3.9))
tf = box.text_frame; tf.word_wrap = True
bullet(tf, "📂 작업 폴더 분리 — 회사·개인 자료 분리", size=13, color=TX, bold=True, first=True)
bullet(tf, "    Claude 는 실행한 폴더만 인식", size=12, color=TX2)
bullet(tf, "🔒 권한 좁히기 — Bash(*) 같은 와일드카드 피하기", size=13, color=TX, bold=True)
bullet(tf, "    예: Bash(python *), Bash(npm *)", size=12, color=TX2)
bullet(tf, "📋 CLAUDE.md 에 보안 규칙 명시", size=13, color=TX, bold=True)
bullet(tf, '    "~/Documents 밖으로 절대 파일 쓰지 말 것"', size=12, color=TX2)

# 우측 위험 명령 카드
add_card(s12, Inches(6.9), Inches(2.3), Inches(6.0), Inches(4.6), accent=BAD)
add_text(s12, Inches(7.2), Inches(2.45), Inches(5.5), Inches(0.5),
         "절대 자동 허용 금지", size=18, bold=True, color=BAD)
box = s12.shapes.add_textbox(Inches(7.2), Inches(2.95), Inches(5.6), Inches(3.9))
tf = box.text_frame; tf.word_wrap = True
bullet(tf, "❌  rm -rf", size=14, color=BAD, bold=True, mono=True, first=True)
bullet(tf, "    파일 영구 삭제 — 복구 불가", size=11, color=TX3)
bullet(tf, "❌  git push --force", size=14, color=BAD, bold=True, mono=True)
bullet(tf, "    팀원의 코드 덮어쓰기", size=11, color=TX3)
bullet(tf, "❌  sudo …", size=14, color=BAD, bold=True, mono=True)
bullet(tf, "    관리자 권한 필요한 시스템 변경", size=11, color=TX3)
bullet(tf, "❌  curl / wget (외부 서버)", size=14, color=BAD, bold=True, mono=True)
bullet(tf, "    매번 직접 확인", size=11, color=TX3)

add_text(s12, Inches(0.6), Inches(7.05), Inches(12), Inches(0.4),
         "💡 위험한 작업은 Hook 으로 자동 차단할 수 있습니다 — 워크북 2.8 절 참고",
         size=12, color=GD)

# ═══════════════════════════════════════════════════════
# Slide 13 — 데이터 · 비용 · 검증
# ═══════════════════════════════════════════════════════
s13 = prs.slides.add_slide(BLANK)
add_bg(s13)
add_header(s13, "TRUST", "데이터 · 비용 · 결과 검증")

cards = [
    ("🔐", "데이터·개인정보", OR, [
        "회사 자료 사용 전 회사 정책 확인",
        "Anthropic 은 입력을 학습에 안 쓰지만 정책 가능성",
        "주민번호·전화번호 등은 Hook 으로 자동 차단",
        "삭제·덮어쓰기 작업 전엔 백업",
    ]),
    ("💰", "비용 관리", RGBColor(0x38, 0xBD, 0xF8), [
        "Pro / Max 구독 — 월 정액으로 충분",
        "API 키 직접 사용은 토큰 단위 과금 (주의)",
        "/usage 로 본인 사용량 수시 확인",
        "Plan 모드 활용 → 토큰 절약",
    ]),
    ("✅", "결과 검증", GD, [
        "AI 산출물은 반드시 사람이 검토",
        "외부 문서·의사결정 자료는 사실 확인 필수",
        "출처 각주, 수치, 톤 한 번씩 더 확인",
        '"AI 가 만들었으니 맞을 것" 가정 금지',
    ]),
]
for i, (icon, title, accent, lines) in enumerate(cards):
    cx = Inches(0.6 + i * 4.13)
    add_card(s13, cx, Inches(2.3), Inches(4), Inches(4.7), accent=accent)
    add_text(s13, cx + Inches(0.3), Inches(2.45), Inches(3.5), Inches(0.6),
             icon + "  " + title, size=15, bold=True, color=accent)
    box = s13.shapes.add_textbox(cx + Inches(0.3), Inches(3.1), Inches(3.6), Inches(3.6))
    tf = box.text_frame; tf.word_wrap = True
    for j, ln in enumerate(lines):
        bullet(tf, "· " + ln, size=12, color=TX, first=(j == 0))

# ═══════════════════════════════════════════════════════
# Slide 14 — 마무리 / 다음 단계
# ═══════════════════════════════════════════════════════
s14 = prs.slides.add_slide(BLANK)
add_bg(s14)
# 큰 오렌지 블록
add_rect(s14, 0, 0, SLIDE_W, Inches(2.3), fill=OR_D)
add_text(s14, Inches(0.6), Inches(0.7), Inches(12), Inches(0.5),
         "NEXT STEP", size=14, bold=True, color=WHITE)
add_text(s14, Inches(0.6), Inches(1.1), Inches(12), Inches(1.0),
         "이제 집에서도 가능합니다", size=36, bold=True, color=WHITE)

# 4 단계 권장
steps = [
    ("1", "오늘", "워크북 PDF · DOCX 를 책처럼 한 번 훑어보기"),
    ("2", "이번 주", "본인 PC 에 Claude Code 설치 + 가벼운 명령 5 개"),
    ("3", "다음 주", "본인 부서 업무 1 개를 선택해서 Skill / Command 만들기"),
    ("4", "한 달 뒤", "Hook · MCP 추가 → 완전 자동화 워크플로우 1 개 완성"),
]
for i, (n, when, what) in enumerate(steps):
    cy = Inches(2.7 + i * 0.85)
    # 번호 원
    circle = s14.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.7), cy + Inches(0.05), Inches(0.65), Inches(0.65))
    circle.fill.solid(); circle.fill.fore_color.rgb = OR; circle.line.fill.background()
    tf = circle.text_frame; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = n; r.font.size = Pt(18); r.font.bold = True; r.font.color.rgb = WHITE; r.font.name = KO_FONT

    add_text(s14, Inches(1.55), cy, Inches(2.8), Inches(0.5),
             when, size=15, bold=True, color=OR)
    add_text(s14, Inches(4.4), cy + Inches(0.05), Inches(8.5), Inches(0.5),
             what, size=14, color=TX)

# 하단 메시지
add_text(s14, Inches(0.6), Inches(6.4), Inches(12), Inches(0.5),
         "막히면 워크북 부록 A 를 다시 펼쳐 보세요.",
         size=14, color=TX2, align=PP_ALIGN.CENTER)
add_text(s14, Inches(0.6), Inches(6.85), Inches(12), Inches(0.5),
         "오늘 배운 한 가지를 일주일 안에 본인 업무에 적용해보는 것이 가장 빠른 실력 향상입니다.",
         size=13, color=GD, align=PP_ALIGN.CENTER)
add_text(s14, Inches(0.6), Inches(7.25), Inches(12), Inches(0.4),
         "미래에셋생명  ·  Vibe Coding Workshop  ·  2026",
         size=10, color=TX3, font=MONO, align=PP_ALIGN.CENTER)

# ═══════════════════════════════════════════════════════
prs.save(str(OUT))
print(f"OK → {OUT} ({OUT.stat().st_size // 1024} KB, {len(prs.slides)} slides)")
